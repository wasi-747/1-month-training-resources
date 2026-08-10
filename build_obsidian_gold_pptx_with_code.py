import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_obsidian_gold_code_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Tokens
    BG_OBSIDIAN = RGBColor(18, 18, 20)       # #121214 Deep Obsidian
    CARD_BG = RGBColor(26, 26, 30)           # #1A1A1E Matte Charcoal
    CODE_BG = RGBColor(15, 15, 18)           # #0F0F12 Deep Terminal Card
    BORDER_SUBTLE = RGBColor(42, 42, 50)     # #2A2A32 Subtle Border
    BORDER_GOLD = RGBColor(245, 158, 11)     # #F59E0B Gold Accent
    
    GOLD_PRIMARY = RGBColor(245, 158, 11)    # #F59E0B Warm Amber-Gold
    GOLD_LIGHT = RGBColor(252, 211, 77)      # #FCD34D Champagne Gold
    WHITE = RGBColor(255, 255, 255)          # #FFFFFF Pure White
    TEXT_BODY = RGBColor(228, 228, 231)      # #E4E4E7 Clean Zinc Body
    TEXT_MUTED = RGBColor(161, 161, 170)     # #A1A1AA Slate Gray
    CODE_GREEN = RGBColor(52, 211, 153)      # #34D399 Terminal Mint Green
    CODE_BLUE = RGBColor(96, 165, 250)       # #60A5FA Electric Sky Blue
    CODE_YELLOW = RGBColor(251, 191, 36)     # #FBBF24 Amber Code

    base_dir = os.path.dirname(os.path.abspath(__file__))
    assets_dir = os.path.join(base_dir, "assets")

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_OBSIDIAN
        bg.line.fill.background()

    def add_header(slide, title_text, category="WEEK 1 TECHNICAL REVIEW"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = GOLD_PRIMARY

        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.6))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(23)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE

        # Divider line
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.733), Inches(0.015))
        div.fill.solid()
        div.fill.fore_color.rgb = BORDER_SUBTLE
        div.line.fill.background()

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_SUBTLE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # ==========================================
    # SLIDE 1: Title Slide (Hero)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1)

    add_card(s1, 0.8, 0.8, 11.733, 5.9)
    tb = s1.shapes.add_textbox(Inches(1.4), Inches(1.3), Inches(10.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "1-MONTH FULL-STACK TRAINING • WEEK 1 REVIEW & DEMO"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = GOLD_PRIMARY
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "Modern Web Foundations &\nAdvanced JavaScript Architecture"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "Comprehensive Architectural Walkthrough with Exact Code Snippets & Line Numbers:\nSemantic HTML5, CSS Grid/Flexbox, Event Loop, Event Delegation, LocalStorage & Node.js Core."
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "🚀 FEATURED LIVE MINI-PROJECT: DevExplorer PRO v2.0 (GitHub Analytics & Bookmark Manager)"
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = GOLD_LIGHT
    p3.space_after = Pt(8)

    p4 = tf.add_paragraph()
    p4.text = "Presenter: Full-Stack Engineering Trainee  |  Date: August 10, 2026"
    p4.font.size = Pt(13)
    p4.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: 4-Day Milestone Roadmap
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2)
    add_header(s2, "Technical Curriculum & 4-Day Milestone Roadmap", "WEEK 1 OVERVIEW")

    days_data = [
        ("DAY 01", "Web Design & Git", [
            ("HTML5 Semantics", "Replaced <div> soup with <main>, <nav>, <header> for SEO & a11y."),
            ("Flexbox & CSS Grid", "1D component alignment vs 2D responsive grid matrices."),
            ("Professional Git", "Feature-branching (git checkout -b) & Conventional Commits.")
        ]),
        ("DAY 02", "Advanced JS & Async", [
            ("ES6+ Architecture", "Block scope (const/let), destructuring, and spread operator."),
            ("Array Pipelines", ".map(), .filter(), and .reduce() 1-pass analytical totals."),
            ("Event Loop Engine", "Call Stack execution, Microtask Queue vs Macrotasks.")
        ]),
        ("DAY 03", "DOM & Mini-Project", [
            ("Event Delegation", "Single parent listener handles dynamic items via closest()."),
            ("LocalStorage Sync", "Client state persistence using JSON serialization."),
            ("DevExplorer PRO", "Parallel API fetching via Promise.all() + live analytics.")
        ]),
        ("DAY 04", "Node.js & FileSystem", [
            ("V8 & Libuv Engine", "Server-side JavaScript runtime with direct OS access."),
            ("fs/promises Module", "Non-blocking, asynchronous disk read & write operations."),
            ("Path Resolution", "Cross-platform path handling via path.join().")
        ])
    ]

    for idx, (day_tag, day_title, items) in enumerate(days_data):
        x = 0.8 + idx * 2.98
        add_card(s2, x, 1.6, 2.8, 5.3)

        tb = s2.shapes.add_textbox(Inches(x + 0.15), Inches(1.75), Inches(2.5), Inches(4.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_d = tf.paragraphs[0]
        p_d.text = day_tag
        p_d.font.size = Pt(11)
        p_d.font.bold = True
        p_d.font.color.rgb = GOLD_PRIMARY
        p_d.space_after = Pt(4)

        p_t = tf.add_paragraph()
        p_t.text = day_title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.space_after = Pt(16)

        for heading, desc in items:
            p_h = tf.add_paragraph()
            p_h.text = "• " + heading
            p_h.font.size = Pt(12)
            p_h.font.bold = True
            p_h.font.color.rgb = GOLD_LIGHT
            p_h.space_after = Pt(2)

            p_desc = tf.add_paragraph()
            p_desc.text = desc
            p_desc.font.size = Pt(10.5)
            p_desc.font.color.rgb = TEXT_MUTED
            p_desc.space_after = Pt(10)

    # ==========================================
    # SLIDE 3: Day 01 - HTML5, Grid/Flexbox & Git Workflow
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3)
    add_header(s3, "Day 01: Modern Web Standards & Team Git Workflows", "DAY 01 ARCHITECTURE")

    # Left: Explanation Panel
    add_card(s3, 0.8, 1.6, 5.7, 5.3)
    tb = s3.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Architectural Decisions & Standards"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

    bullets_s3 = [
        ("Semantic HTML5 Hierarchy (index.html: L30-L240)", "Utilized <header>, <nav>, <main>, <aside>, <article>, and <footer>. Replaced generic div soup to maximize Google SEO and screen reader accessibility."),
        ("1D Flexbox vs 2D CSS Grid (styles.css)", "Flexbox aligns 1D axes (.header-container, .quick-tags) with justify-content: space-between. CSS Grid powers 2D dashboard layouts."),
        ("The Golden Responsive Grid Formula", "grid-template-columns: repeat(auto-fit, minmax(280px, 1fr)) creates automatic multi-column responsiveness without writing CSS media queries!"),
        ("Professional Feature-Branching Git Workflow", "All changes isolated in feature branches (git checkout -b) with Conventional Commits (feat:, fix:, docs:) before merging via PRs.")
    ]
    for h, d in bullets_s3:
        ph = tf.add_paragraph()
        ph.text = "✔  " + h
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = GOLD_LIGHT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    # Right: Code Snippet Card
    add_card(s3, 6.8, 1.6, 5.733, 5.3, bg_color=CODE_BG)
    tb_c = s3.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0

    pc_title = tf_c.paragraphs[0]
    pc_title.text = "📄 Production Snippet: HTML5, CSS Grid & Git"
    pc_title.font.size = Pt(13)
    pc_title.font.bold = True
    pc_title.font.color.rgb = GOLD_PRIMARY
    pc_title.space_after = Pt(10)

    code_s3 = (
        "/* 1. CSS GRID: Golden Auto-Responsive Formula */\n"
        ".repos-grid {\n"
        "  display: grid;\n"
        "  grid-template-columns: repeat(auto-fit, minmax(280px, 1fr));\n"
        "  gap: 1.5rem;\n"
        "}\n\n"
        "/* 2. FLEXBOX: 1D Single-Axis Navbar */\n"
        ".header-container {\n"
        "  display: flex;\n"
        "  justify-content: space-between; /* Space logo & tabs */\n"
        "  align-items: center;            /* Vertical center */\n"
        "}\n\n"
        "# 3. Professional Git Team Workflow\n"
        "$ git checkout -b feature/dashboard-ui\n"
        "$ git add index.html styles.css\n"
        "$ git commit -m \"feat: build semantic dashboard\"\n"
        "$ git push -u origin feature/dashboard-ui"
    )
    pc_body = tf_c.add_paragraph()
    pc_body.text = code_s3
    pc_body.font.name = "Consolas"
    pc_body.font.size = Pt(10)
    pc_body.font.color.rgb = CODE_GREEN

    # ==========================================
    # SLIDE 4: Day 02 - ES6+ Scope & The Event Loop
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4)
    add_header(s4, "Day 02: Advanced JS Scope & Event Loop Engine", "DAY 02 CORE ENGINE")

    # Left: Explanation Panel
    add_card(s4, 0.8, 1.6, 5.7, 5.3)
    tb = s4.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Memory, Scope & Engine Concurrency"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

    bullets_s4 = [
        ("Block Scope & TDZ (app.js: L162, L273)", "const / let are confined strictly within enclosing {}. Accessing them before initialization throws ReferenceError (TDZ), eliminating var hoisting bugs."),
        ("Single-Threaded Call Stack LIFO (app.js: L206)", "V8 executes synchronous functions line-by-line (Last-In, First-Out). Pushes renderProfileCard(), executes, pops it off, then pushes computeAnalytics()."),
        ("Microtask VIP Queue (app.js: L181, L197)", "Promise.all() and await res.json() resolve in the Microtask Queue with VIP priority, draining immediately when Call Stack clears before any Macrotask."),
        ("Macrotask Queue (app.js: L497)", "setTimeout(() => toast.remove(), 3000) enters Macrotask queue, executing only after all pending microtasks are complete.")
    ]
    for h, d in bullets_s4:
        ph = tf.add_paragraph()
        ph.text = "⚡  " + h
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = GOLD_LIGHT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    # Right: Code Snippet Card
    add_card(s4, 6.8, 1.6, 5.733, 5.3, bg_color=CODE_BG)
    tb_c = s4.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0

    pc_title = tf_c.paragraphs[0]
    pc_title.text = "📄 app.js: Engine Concurrency & Scope"
    pc_title.font.size = Pt(13)
    pc_title.font.bold = True
    pc_title.font.color.rgb = GOLD_PRIMARY
    pc_title.space_after = Pt(8)

    code_s4 = (
        "// 1. Block Scope & TDZ (app.js: L162)\n"
        "const cleanUsername = username.toLowerCase();\n"
        "// Accessing cleanUsername before L162 = TDZ ReferenceError\n\n"
        "// 2. Microtask VIP Queue (app.js: L181)\n"
        "const [userRes, reposRes] = await Promise.all([\n"
        "  fetch(`https://api.github.com/users/${cleanUsername}`),\n"
        "  fetch(`https://api.github.com/users/${cleanUsername}/repos`)\n"
        "]); // Resolves with VIP priority in Microtask Queue!\n\n"
        "// 3. Single-Threaded Call Stack LIFO (app.js: L206)\n"
        "this.renderProfileCard(userData, false); // Stack Push -> Pop\n"
        "this.computeAnalytics(reposData);        // Stack Push -> Pop\n\n"
        "// 4. Macrotask Queue (app.js: L497)\n"
        "setTimeout(() => toast.remove(), 3000);  // Runs after Microtasks"
    )
    pc_body = tf_c.add_paragraph()
    pc_body.text = code_s4
    pc_body.font.name = "Consolas"
    pc_body.font.size = Pt(9.5)
    pc_body.font.color.rgb = CODE_YELLOW

    # ==========================================
    # SLIDE 5: Mini-Project Architecture - Async & Promise.all()
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5)
    add_header(s5, "Mini-Project: Parallel API Fetching & Caching", "APP ARCHITECTURE")

    # Left: Explanation Panel
    add_card(s5, 0.8, 1.6, 5.7, 5.3)
    tb = s5.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Parallel Network Engine & In-Memory Cache"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

    bullets_s5 = [
        ("In-Memory Map Cache (app.js: L23, L166)", "ES6 Private Field #searchCache = new Map(). Checks if user was already fetched; re-renders in 0ms without hitting GitHub rate limits."),
        ("Promise.all Concurrent Fetching (app.js: L181)", "Dispatches User Profile and Repositories endpoints simultaneously in parallel, cutting network latency by 50% compared to sequential awaits."),
        ("Stream Decoding (app.js: L197)", "await userRes.json() decodes the raw incoming HTTP stream into live JavaScript objects."),
        ("Try / Catch Error Boundaries (app.js: L211)", "Catches network dropouts and displays friendly UI alerts without crashing the runtime application.")
    ]
    for h, d in bullets_s5:
        ph = tf.add_paragraph()
        ph.text = "🚀  " + h
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = GOLD_LIGHT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    # Right: Code Snippet Card
    add_card(s5, 6.8, 1.6, 5.733, 5.3, bg_color=CODE_BG)
    tb_c = s5.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0

    pc_title = tf_c.paragraphs[0]
    pc_title.text = "📄 app.js: Lines 166–203 (Parallel Fetch & Cache)"
    pc_title.font.size = Pt(13)
    pc_title.font.bold = True
    pc_title.font.color.rgb = GOLD_PRIMARY
    pc_title.space_after = Pt(10)

    code_s5 = (
        "// 1. Check In-Memory Map Cache (0ms Instant Load)\n"
        "if (this.#searchCache.has(cleanUsername)) {\n"
        "  const cached = this.#searchCache.get(cleanUsername);\n"
        "  this.renderProfileCard(cached.user, true);\n"
        "  return;\n"
        "}\n\n"
        "// 2. Parallel API Dispatch via Promise.all()\n"
        "const [userRes, reposRes] = await Promise.all([\n"
        "  fetch(`https://api.github.com/users/${username}`),\n"
        "  fetch(`https://api.github.com/users/${username}/repos?per_page=100`)\n"
        "]);\n\n"
        "// 3. Stream Parsing & State Persistence\n"
        "const userData = await userRes.json();\n"
        "const reposData = reposRes.ok ? await reposRes.json() : [];\n"
        "this.#searchCache.set(cleanUsername, { user: userData, repos: reposData });"
    )
    pc_body = tf_c.add_paragraph()
    pc_body.text = code_s5
    pc_body.font.name = "Consolas"
    pc_body.font.size = Pt(9.5)
    pc_body.font.color.rgb = CODE_GREEN

    # ==========================================
    # SLIDE 6: Mini-Project Analytics Engine (.reduce)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6)
    add_header(s6, "Mini-Project: Single-Pass Analytics Engine", "DATA PIPELINES")

    # Left: Explanation Panel
    add_card(s6, 0.8, 1.6, 5.7, 5.3)
    tb = s6.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Higher-Order Analytical Aggregation"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

    bullets_s6 = [
        ("Total Stars in 1-Pass (app.js: L228)", "repos.reduce((acc, repo) => acc + repo.stargazers_count, 0) aggregates total star counts across 100+ repositories in O(N) linear time."),
        ("Total Forks Aggregation (app.js: L232)", "repos.reduce((acc, repo) => acc + repo.forks_count, 0) accumulates community fork metrics."),
        ("Top Language Frequency Map (app.js: L237)", ".reduce() constructs a dynamic frequency dictionary: { JavaScript: 15, TypeScript: 8, Python: 3 } to determine developer specialty."),
        ("Immutability & Spread (app.js: L259, L273)", "const languages = [...new Set(...)]; and let result = [...this.currentRepos] guarantee original API state is never mutated.")
    ]
    for h, d in bullets_s6:
        ph = tf.add_paragraph()
        ph.text = "📊  " + h
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = GOLD_LIGHT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    # Right: Code Snippet Card
    add_card(s6, 6.8, 1.6, 5.733, 5.3, bg_color=CODE_BG)
    tb_c = s6.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0

    pc_title = tf_c.paragraphs[0]
    pc_title.text = "📄 app.js: Lines 228–248 (computeAnalytics)"
    pc_title.font.size = Pt(13)
    pc_title.font.bold = True
    pc_title.font.color.rgb = GOLD_PRIMARY
    pc_title.space_after = Pt(10)

    code_s6 = (
        "// 1. Total Stars Calculation (.reduce single-pass)\n"
        "const totalStars = repos.reduce((acc, repo) => \n"
        "  acc + (repo.stargazers_count || 0), 0\n"
        ");\n"
        "this.metricTotalStars.textContent = totalStars.toLocaleString();\n\n"
        "// 2. Language Frequency Map Generator (.reduce)\n"
        "const langCounts = repos.reduce((acc, repo) => {\n"
        "  if (repo.language) {\n"
        "    acc[repo.language] = (acc[repo.language] || 0) + 1;\n"
        "  }\n"
        "  return acc;\n"
        "}, {}); // Output: { JavaScript: 15, Python: 4 }\n\n"
        "// 3. Extract Dominant Language\n"
        "let topLang = 'JavaScript', maxCount = 0;\n"
        "for (const [lang, count] of Object.entries(langCounts)) {\n"
        "  if (count > maxCount) { maxCount = count; topLang = lang; }\n"
        "}"
    )
    pc_body = tf_c.add_paragraph()
    pc_body.text = code_s6
    pc_body.font.name = "Consolas"
    pc_body.font.size = Pt(9.5)
    pc_body.font.color.rgb = CODE_YELLOW

    # ==========================================
    # SLIDE 7: Mini-Project Event Delegation & LocalStorage
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7)
    add_header(s7, "Mini-Project: Event Delegation & LocalStorage", "DOM & STORAGE")

    # Left: Explanation Panel
    add_card(s7, 0.8, 1.6, 5.7, 5.3)
    tb = s7.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Memory Efficiency & Client Persistence"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

    bullets_s7 = [
        ("Single Container Listener (app.js: L102, L141)", "Instead of binding 50+ listeners, attached 1 listener to #bookmarks-grid and #quick-tags-container. Prevents memory leaks."),
        ("Event Bubbling & closest() (app.js: L104, L142)", "Clicks bubble upwards; e.target.closest('.tag-btn') reliably extracts dataset.user even if user clicked inner icons or spans."),
        ("JSON Serialization Protocol (app.js: L505, L517)", "localStorage only stores strings. Used JSON.stringify() to serialize array state, and JSON.parse() on load."),
        ("State Sync (.some & Badges) (app.js: L406)", "this.bookmarks.some() syncs Star bookmark toggle button and updates navbar badge counts automatically.")
    ]
    for h, d in bullets_s7:
        ph = tf.add_paragraph()
        ph.text = "🎯  " + h
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = GOLD_LIGHT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    # Right: Code Snippet Card
    add_card(s7, 6.8, 1.6, 5.733, 5.3, bg_color=CODE_BG)
    tb_c = s7.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0

    pc_title = tf_c.paragraphs[0]
    pc_title.text = "📄 app.js: Lines 102–145, 505–523"
    pc_title.font.size = Pt(13)
    pc_title.font.bold = True
    pc_title.font.color.rgb = GOLD_PRIMARY
    pc_title.space_after = Pt(10)

    code_s7 = (
        "// 1. Event Delegation on Parent Grid (app.js: L141)\n"
        "this.bookmarksGrid.addEventListener('click', (e) => {\n"
        "  const btnInspect = e.target.closest('.btn-inspect-dev');\n"
        "  const btnRemove = e.target.closest('.btn-remove-bookmark');\n"
        "  if (btnInspect) this.fetchDeveloperProfile(btnInspect.dataset.user);\n"
        "  if (btnRemove) this.removeBookmark(btnRemove.dataset.user);\n"
        "});\n\n"
        "// 2. LocalStorage Read Protocol (app.js: L505)\n"
        "#loadBookmarksFromStorage() {\n"
        "  const data = localStorage.getItem(this.#STORAGE_KEY);\n"
        "  return data ? JSON.parse(data) : []; // String -> Array\n"
        "}\n\n"
        "// 3. LocalStorage Write Protocol (app.js: L517)\n"
        "#saveBookmarksToStorage() {\n"
        "  localStorage.setItem(this.#STORAGE_KEY, JSON.stringify(this.bookmarks));\n"
        "}"
    )
    pc_body = tf_c.add_paragraph()
    pc_body.text = code_s7
    pc_body.font.name = "Consolas"
    pc_body.font.size = Pt(9.5)
    pc_body.font.color.rgb = CODE_BLUE

    # ==========================================
    # SLIDE 8: Day 04 - Node.js V8 & fs/promises
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8)
    add_header(s8, "Day 04: Node.js V8 Runtime & Asynchronous FileSystem", "DAY 04 BACKEND")

    # Left: Explanation Panel
    add_card(s8, 0.8, 1.6, 5.7, 5.3)
    tb = s8.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Server-Side V8 & Non-Blocking FileSystem"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

    bullets_s8 = [
        ("V8 Runtime + Libuv C++ Layer", "Chrome V8 translates JS to machine code while multi-threaded Libuv executes background disk I/O without blocking server traffic."),
        ("Non-Blocking fs/promises Module", "Replaced synchronous fs.readFileSync with await fs.readFile / fs.writeFile. Keeps main Event Loop 100% responsive."),
        ("Cross-Platform path.join(__dirname, ...)", "Normalizes path separators (Windows backslash vs POSIX forward slash) preventing production deployment path crashes."),
        ("Backend Analytics Logging Pipeline", "Automated Node.js script reads, aggregates, and persists DevExplorer telemetry snapshots to disk asynchronously.")
    ]
    for h, d in bullets_s8:
        ph = tf.add_paragraph()
        ph.text = "🟢  " + h
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = GOLD_LIGHT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    # Right: Code Snippet Card
    add_card(s8, 6.8, 1.6, 5.733, 5.3, bg_color=CODE_BG)
    tb_c = s8.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0

    pc_title = tf_c.paragraphs[0]
    pc_title.text = "📄 Node.js: practice-snippets.js (fs/promises)"
    pc_title.font.size = Pt(13)
    pc_title.font.bold = True
    pc_title.font.color.rgb = GOLD_PRIMARY
    pc_title.space_after = Pt(10)

    code_s8 = (
        "const fs = require('fs/promises');\n"
        "const path = require('path');\n\n"
        "async function saveAnalyticsLog(analyticsData) {\n"
        "  try {\n"
        "    // 1. Cross-Platform Safe Path\n"
        "    const logPath = path.join(__dirname, 'logs', 'analytics.json');\n\n"
        "    // 2. Non-Blocking Asynchronous Disk Write\n"
        "    await fs.writeFile(logPath, JSON.stringify(analyticsData, null, 2));\n"
        "    console.log('✔ Non-blocking write completed!');\n\n"
        "    // 3. Asynchronous Non-Blocking Disk Read\n"
        "    const raw = await fs.readFile(logPath, 'utf-8');\n"
        "    return JSON.parse(raw);\n"
        "  } catch (err) {\n"
        "    console.error('Disk I/O Error:', err.message);\n"
        "  }\n"
        "}"
    )
    pc_body = tf_c.add_paragraph()
    pc_body.text = code_s8
    pc_body.font.name = "Consolas"
    pc_body.font.size = Pt(9.5)
    pc_body.font.color.rgb = CODE_GREEN

    # ==========================================
    # SLIDE 9: Live Demo Cheat Sheet & Walkthrough
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9)
    add_header(s9, "Live Demonstration Walkthrough & Cheat Sheet", "LIVE DEMO")

    demo_steps = [
        ("Step 1: Search Trigger", "Type 'torvalds' or click popular developer tags", "app.js: L94", "Promise.all concurrent API fetch executes; Network tab shows parallel user & repo responses."),
        ("Step 2: Analytics Cards", "Inspect Total Stars, Forks & Language badges", "app.js: L228", ".reduce() aggregates 100+ repos into single total numbers in 1 linear pass."),
        ("Step 3: Live Filters", "Type keywords in filter box & select language", "app.js: L273", "Spread operator maintains immutability while .filter() updates grid without reloading."),
        ("Step 4: Bookmark & Storage", "Click 'Save to Bookmarks' & inspect Tab", "app.js: L369", "JSON.stringify saves to localStorage; badge count updates; toast animates via setTimeout."),
        ("Step 5: Event Delegation", "Click 'Inspect' or 'Remove' in Bookmarks tab", "app.js: L141", "1 single parent listener intercepts clicks via event.target.closest(). Zero memory leaks!"),
        ("Step 6: Cache Verification", "Search 'torvalds' again in Explorer tab", "app.js: L166", "In-Memory Map cache loads profile in 0ms instant speed; 'Cached ⚡' indicator activates.")
    ]

    for idx, (title, action, line_ref, tech_detail) in enumerate(demo_steps):
        col = idx % 2
        row = idx // 2
        x = 0.8 + col * 5.95
        y = 1.6 + row * 1.75

        add_card(s9, x, y, 5.75, 1.6)

        tb = s9.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(5.45), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = title + "  (" + line_ref + ")"
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = GOLD_LIGHT
        p1.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = "Action: " + action
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.space_after = Pt(2)

        p3 = tf.add_paragraph()
        p3.text = "Under the Hood: " + tech_detail
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 10: Summary & Week 2 Transition
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10)
    add_header(s10, "Summary, Key Learnings & Next Steps", "LOOKING FORWARD")

    add_card(s10, 0.8, 1.6, 5.7, 5.3)
    tb = s10.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Week 1 Core Learnings"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

    learnings = [
        ("Production Web Standards", "Adopted Semantic HTML5, 1D Flexbox, and 2D Auto-responsive CSS Grid."),
        ("Advanced JS & Event Loop", "Learned block scoping, immutability, Microtasks vs Macrotasks, and Promise concurrency."),
        ("Event Delegation Architecture", "Eliminated memory leaks using single parent container listeners with closest()."),
        ("Server-Side Node.js Runtime", "Understood V8 engine, Libuv thread pool, and non-blocking fs/promises disk I/O.")
    ]
    for h, d in learnings:
        ph = tf.add_paragraph()
        ph.text = "✔  " + h
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = GOLD_LIGHT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    # Right: Week 2 Roadmap Card
    add_card(s10, 6.8, 1.6, 5.733, 5.3)
    tb = s10.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Week 2: React Deep Dive Roadmap"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = GOLD_PRIMARY
    p.space_after = Pt(14)

    w2_plan = [
        ("Thinking in React & JSX (Day 01)", "Component hierarchies, JSX compilation rules, and Props vs State architecture."),
        ("Interactive State & Hooks (Day 02)", "useState mechanics, immutable state updates, and building a modular Tic-Tac-Toe game."),
        ("Effects & Custom Hooks (Day 03)", "useEffect lifecycle management, dependency arrays, cleanup functions, and custom API hooks."),
        ("Context API & Mini-Project (Day 04)", "Global state management without prop-drilling, culminating in a Multi-Theme Task Tracker.")
    ]
    for h, d in w2_plan:
        ph = tf.add_paragraph()
        ph.text = "🚀  " + h
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = GOLD_LIGHT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    # Save with lock fallback
    out_path = os.path.join(base_dir, "Week_1_Presentation_Obsidian_Gold.pptx")
    try:
        prs.save(out_path)
        print(f"SUCCESS: Obsidian Gold Deck saved to {out_path}")
    except PermissionError:
        fallback_path = os.path.join(base_dir, "Week_1_Presentation_Obsidian_Gold_v2.pptx")
        prs.save(fallback_path)
        print(f"SUCCESS: File locked in PowerPoint. Saved updated deck to {fallback_path}")

if __name__ == "__main__":
    create_obsidian_gold_code_deck()
