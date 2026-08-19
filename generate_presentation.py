# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- 🎨 Color Palette (Modern Dark Slate Executive Theme) ---
DARK_SLATE = RGBColor(15, 23, 42)    # #0F172A
SLATE_800 = RGBColor(30, 41, 59)     # #1E293B
SLATE_850 = RGBColor(24, 33, 47)     # #18212F
SKY_CYAN = RGBColor(56, 189, 248)    # #38BDF8
EMERALD = RGBColor(52, 211, 153)     # #34D399
AMBER = RGBColor(251, 191, 36)       # #FBBF24
ROSE = RGBColor(251, 113, 133)       # #FB7185
INDIGO = RGBColor(129, 140, 248)     # #818CF8
WHITE = RGBColor(255, 255, 255)      # #FFFFFF
LIGHT_SLATE = RGBColor(226, 232, 240)# #E2E8F0
MUTED = RGBColor(148, 163, 184)      # #94A3B8

# --- 🛠️ Helper Functions ---
def create_blank_slide(prs):
    """Creates a blank slide with the Dark Slate background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 is a blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_SLATE
    return slide

def add_text(slide, text, left, top, width, height, font_size, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Adds a simple text box to the slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = 'Segoe UI'
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox

def add_card(slide, title, bullets, left, top, width, height, accent_color=SKY_CYAN):
    """Draws a styled rounded rectangle card with a title and bullet points."""
    # Create Rounded Rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SLATE_800
    shape.line.color.rgb = SLATE_850 # Subtle border
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    
    # Title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title + "\n"
    run.font.name = 'Segoe UI'
    run.font.size = Pt(20)
    run.font.color.rgb = accent_color
    run.font.bold = True
    
    # Bullets
    for bullet in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.name = 'Segoe UI'
        run.font.size = Pt(16)
        run.font.color.rgb = LIGHT_SLATE

def add_slide_title(slide, title_text, subtitle_text=None):
    """Standardized title header for content slides."""
    add_text(slide, title_text, 0.5, 0.3, 12.33, 0.8, 32, WHITE, True)
    if subtitle_text:
        add_text(slide, subtitle_text, 0.5, 0.9, 12.33, 0.5, 18, MUTED)

# --- 🚀 Build Presentation ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---------------------------------------------------------
# SLIDE 1: Title Slide
# ---------------------------------------------------------
slide1 = create_blank_slide(prs)
add_text(slide1, "WEEK 2 EVALUATION PRESENTATION", 1, 2.2, 11.33, 1.0, 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide1, "React Deep Dive, MongoDB & Next.js Architecture", 1, 3.2, 11.33, 0.8, 24, SKY_CYAN, False, PP_ALIGN.CENTER)

# Badges (Simulated using text background or just colored text)
badges = "[⚛️ React 18]   [⚡ Custom Hooks]   [🍃 MongoDB]   [🛡️ Mongoose]   [▲ Next.js]   [⚡ Vite]"
add_text(slide1, badges, 1, 4.2, 11.33, 0.5, 16, EMERALD, True, PP_ALIGN.CENTER)

add_text(slide1, "Presenter: Wasiur Rahman Sakib | Software Engineer Intern", 1, 5.5, 11.33, 0.5, 18, MUTED, False, PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 2: Executive Overview (4 Core Modules)
# ---------------------------------------------------------
slide2 = create_blank_slide(prs)
add_slide_title(slide2, "Executive Overview", "Week 2: 4 Core Modules")

# 2x2 Grid
w, h = 5.8, 2.4
add_card(slide2, "React Paradigm", ["Declarative UI", "Virtual DOM Diffing", "Unidirectional Data Flow"], 0.6, 1.8, w, h, SKY_CYAN)
add_card(slide2, "State & Hooks", ["Immutability Principles", "useEffect Lifecycle", "Custom Hooks"], 6.8, 1.8, w, h, EMERALD)
add_card(slide2, "Database Layer", ["MongoDB BSON", "Mongoose Schemas", "Validation & Aggregations"], 0.6, 4.5, w, h, AMBER)
add_card(slide2, "Full-Stack Bridge", ["Next.js App Router", "Server vs Client Components", "API Routes"], 6.8, 4.5, w, h, INDIGO)

# ---------------------------------------------------------
# SLIDE 3: React Core & Immutability
# ---------------------------------------------------------
slide3 = create_blank_slide(prs)
add_slide_title(slide3, "React Core & Immutability")

w, h = 3.8, 3.5
add_card(slide3, "Uni-directional Flow", ["Props Down", "Events Up", "Single Source of Truth"], 0.6, 2.0, w, h, SKY_CYAN)
add_card(slide3, "Virtual DOM", ["In-Memory Tree", "Reconciliation Diffing", "Batched DOM Updates"], 4.7, 2.0, w, h, EMERALD)
add_card(slide3, "State Immutability", ["Reference Equality (prev !== next)", "Spread [...prev]", "Time-Travel Snapshots"], 8.8, 2.0, w, h, ROSE)

# ---------------------------------------------------------
# SLIDE 4: Lifecycle & Custom Hooks Architecture
# ---------------------------------------------------------
slide4 = create_blank_slide(prs)
add_slide_title(slide4, "Lifecycle & Custom Hooks Architecture")

add_card(slide4, "useEffect Sync", ["External System Synchronization", "Controlled Dependency Arrays"], 0.6, 2.0, w, h, AMBER)
add_card(slide4, "Cleanup Phase", ["Abort Pending Requests", "Clear Listeners", "Prevent Memory Leaks"], 4.7, 2.0, w, h, ROSE)
add_card(slide4, "Custom Hooks", ["Reusable Logic", "useFetch (Async Lifecycle)", "useToggle (Boolean UI)"], 8.8, 2.0, w, h, SKY_CYAN)

# ---------------------------------------------------------
# SLIDE 5: MongoDB, Mongoose & Next.js Foundations
# ---------------------------------------------------------
slide5 = create_blank_slide(prs)
add_slide_title(slide5, "MongoDB, Mongoose & Next.js Foundations")

add_card(slide5, "MongoDB BSON", ["Flexible Documents", "Embedded Sub-docs", "Aggregation Pipelines"], 0.6, 2.0, w, h, EMERALD)
add_card(slide5, "Mongoose ODM", ["Schema Validation", "Type Casting", "Pre/Post Middleware Hooks"], 4.7, 2.0, w, h, AMBER)
add_card(slide5, "Next.js Paradigm", ["Server Components (Zero JS)", "Client Interactivity", "API Endpoints"], 8.8, 2.0, w, h, INDIGO)

# ---------------------------------------------------------
# SLIDE 6: Technical Challenges & Engineering Solutions
# ---------------------------------------------------------
slide6 = create_blank_slide(prs)
add_slide_title(slide6, "Technical Challenges & Engineering Solutions")

def add_solution_row(slide, top, problem, solution):
    # Problem Badge
    shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(top), Inches(5.8), Inches(1.2))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = SLATE_850
    tf1 = shape1.text_frame
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = f"⚠️ {problem}"
    run1.font.name, run1.font.size, run1.font.color.rgb = 'Segoe UI', Pt(18), ROSE
    
    # Solution Badge
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(top), Inches(5.8), Inches(1.2))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = SLATE_800
    tf2 = shape2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = f"💡 {solution}"
    run2.font.name, run2.font.size, run2.font.color.rgb = 'Segoe UI', Pt(18), EMERALD

add_solution_row(slide6, 1.8, "Async Race Conditions", "AbortController Signal & useEffect Cleanup")
add_solution_row(slide6, 3.3, "State Array Mutations", "Spread Operator [...prev] & Pure Array Maps")
add_solution_row(slide6, 4.8, "DB Connection Spawning", "Global Singleton Cached Connection")

# ---------------------------------------------------------
# SLIDE 7: Practical Demonstration (Live Projects)
# ---------------------------------------------------------
slide7 = create_blank_slide(prs)
add_slide_title(slide7, "Practical Demonstration", "Live Projects & Implementations")

add_card(slide7, "Tic-Tac-Toe Arena", ["Immutable Board Snapshots", "Time-Travel Move History", "Pure Win Calculation"], 1.0, 1.8, 5.2, 3.5, SKY_CYAN)
add_card(slide7, "Custom Hooks Suite", ["useFetch Async States", "useToggle UI Triggers", "Decoupled Logic"], 7.0, 1.8, 5.2, 3.5, EMERALD)

# Bottom Banner with URL
banner = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.7), Inches(11.2), Inches(1.0))
banner.fill.solid()
banner.fill.fore_color.rgb = INDIGO
banner.line.fill.background()
tf = banner.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "🚀 Live Deployed App URL: "
run.font.name, run.font.size, run.font.color.rgb, run.font.bold = 'Segoe UI', Pt(18), WHITE, True
url_run = p.add_run()
url_run.text = "Click Here to View Live (GitHub Pages)"
url_run.font.name, url_run.font.size, url_run.font.color.rgb, url_run.font.bold = 'Segoe UI', Pt(18), WHITE, True
url_run.hyperlink.address = "https://wasi-747.github.io/1-month-training-resources/Week-2-React-Deep-Dive/react-playground/dist/"

# ---------------------------------------------------------
# SLIDE 8: Engineering Principles & Best Practices
# ---------------------------------------------------------
slide8 = create_blank_slide(prs)
add_slide_title(slide8, "Engineering Principles & Best Practices")

w, h = 5.8, 2.0
add_card(slide8, "01. 🎯 Declarative UI", ["State-driven rendering over manual DOM mutations"], 0.6, 1.8, w, h, WHITE)
add_card(slide8, "02. 🔗 Single Source of Truth", ["Lift state up to eliminate desync bugs"], 6.8, 1.8, w, h, WHITE)
add_card(slide8, "03. 🛡️ Defensive UI Flow", ["Loading skeletons, error boundaries, fallback states"], 0.6, 4.2, w, h, WHITE)
add_card(slide8, "04. 📦 Modular Logic", ["Clean separation of UI views and custom hooks"], 6.8, 4.2, w, h, WHITE)

# ---------------------------------------------------------
# SLIDE 9: Academic Leave & Week 4 Roadmap
# ---------------------------------------------------------
slide9 = create_blank_slide(prs)
add_slide_title(slide9, "Roadmap: Looking Ahead")

add_card(slide9, "Week 3: Academic Leave", ["University Midterm Examinations", "Mental consolidation of Web & React"], 1.0, 2.2, 5.2, 3.5, ROSE)
add_card(slide9, "Week 4: React Native Mobile Dev", ["Cross-Platform iOS & Android", "Core Native Components", "React Navigation", "Hardware Device APIs", "Full-Stack Backend Integration"], 7.0, 2.2, 5.2, 3.5, SKY_CYAN)

# ---------------------------------------------------------
# SLIDE 10: Conclusion & Q&A
# ---------------------------------------------------------
slide10 = create_blank_slide(prs)

add_text(slide10, "THANK YOU!", 1, 2.2, 11.33, 1.0, 54, EMERALD, True, PP_ALIGN.CENTER)
add_text(slide10, "Questions & Live Project Walkthrough", 1, 3.5, 11.33, 0.8, 24, LIGHT_SLATE, False, PP_ALIGN.CENTER)

# Hyperlink Textbox
txBox = slide10.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(0.8))
p = txBox.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Live Game URL 🔗"
run.font.name, run.font.size, run.font.color.rgb, run.font.bold = 'Segoe UI', Pt(20), SKY_CYAN, True
run.hyperlink.address = "https://wasi-747.github.io/1-month-training-resources/Week-2-React-Deep-Dive/react-playground/dist/"

add_text(slide10, "Presenter: Wasiur Rahman Sakib | Software Engineer Intern", 1, 5.5, 11.33, 0.5, 18, MUTED, False, PP_ALIGN.CENTER)


# --- 💾 Save the Presentation ---
output_file = "Week_2_Evaluation_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation successfully generated and saved as '{output_file}'!")