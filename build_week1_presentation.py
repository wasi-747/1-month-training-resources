import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    # 16:9 Widescreen standard (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Premium Dark Theme Color Palette
    BG_COLOR = RGBColor(11, 15, 25)        # #0B0F19 Deep Night Slate
    SURFACE = RGBColor(22, 30, 49)          # #161E31 Surface Card
    SURFACE_LIGHT = RGBColor(30, 41, 59)    # #1E293B Card Inner
    INDIGO = RGBColor(99, 102, 241)         # #6366F1 Accent Indigo
    CYAN = RGBColor(6, 182, 212)           # #06B6D4 Accent Cyan
    AMBER = RGBColor(245, 158, 11)          # #F59E0B Accent Amber
    GREEN = RGBColor(16, 185, 129)          # #10B981 Emerald Green
    WHITE = RGBColor(248, 250, 252)         # #F8FAFC Text Main
    MUTED = RGBColor(148, 163, 184)         # #94A3B8 Text Secondary
    BORDER_SUBTLE = RGBColor(40, 53, 80)

    base_dir = os.path.dirname(os.path.abspath(__file__))
    assets_dir = os.path.join(base_dir, "assets")

    def apply_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

    def add_card(slide, left, top, width, height, bg_rgb=SURFACE, border_rgb=INDIGO, border_w=1.5):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_rgb
        card.line.color.rgb = border_rgb
        card.line.width = Pt(border_w)
        return card

    def add_header(slide, title_text, category="WEEK 1 • ENGINEERING MASTERY"):
        # Top Category Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(3.6), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = INDIGO
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.word_wrap = True
        pb = tf_b.paragraphs[0]
        pb.text = category
        pb.font.size = Pt(10.5)
        pb.font.bold = True
        pb.font.color.rgb = WHITE
        pb.alignment = PP_ALIGN.CENTER

        # Main Slide Title
        tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_bg(s1)

    add_card(s1, 0.8, 0.9, 11.733, 5.7, bg_rgb=SURFACE, border_rgb=INDIGO, border_w=2)
    
    b1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.3), Inches(3.8), Inches(0.45))
    b1.fill.solid()
    b1.fill.fore_color.rgb = INDIGO
    b1.line.fill.background()
    p_b1 = b1.text_frame.paragraphs[0]
    p_b1.text = "1-MONTH TRAINING • PRESENTATION #1"
    p_b1.font.size = Pt(11)
    p_b1.font.bold = True
    p_b1.font.color.rgb = WHITE
    p_b1.alignment = PP_ALIGN.CENTER

    tbox1 = s1.shapes.add_textbox(Inches(1.3), Inches(1.9), Inches(10.8), Inches(2.6))
    tf1 = tbox1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "Modern Web Foundations & Advanced JavaScript"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Semantic Architecture • Professional Git • Async Event Loop • DOM Delegation • Node.js Core"
    p1_sub.font.size = Pt(16.5)
    p1_sub.font.color.rgb = CYAN
    p1_sub.space_before = Pt(10)

    # Footer Card with Live Demo Callout
    demo_callout = add_card(s1, 1.3, 4.8, 10.733, 1.3, bg_rgb=SURFACE_LIGHT, border_rgb=AMBER, border_w=1.5)
    tb_demo = s1.shapes.add_textbox(Inches(1.5), Inches(4.9), Inches(10.3), Inches(1.1))
    tfd = tb_demo.text_frame
    tfd.word_wrap = True
    pd1 = tfd.paragraphs[0]
    pd1.text = "🚀 FEATURED LIVE PROJECT DEMONSTRATION:"
    pd1.font.size = Pt(12)
    pd1.font.bold = True
    pd1.font.color.rgb = AMBER

    pd2 = tfd.add_paragraph()
    pd2.text = "DevExplorer PRO v2.0 — GitHub Analytics Dashboard with Parallel API Fetching & Local Storage"
    pd2.font.size = Pt(14)
    pd2.font.bold = True
    pd2.font.color.rgb = WHITE
    pd2.space_before = Pt(4)

    # ==========================================
    # SLIDE 2: AGENDA & WEEKLY ROADMAP
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_bg(s2)
    add_header(s2, "Week 1 Technical Curriculum & Milestone Roadmap")

    days_data = [
        ("Day 01", "Web Design & Git", "Semantic HTML5, CSS Grid vs Flexbox, Tailwind CSS, Feature-branching & Conventional Commits.", INDIGO),
        ("Day 02", "Advanced JS & Async", "ES6+ array pipelines (.map, .filter, .reduce), Event Loop, Call Stack, Microtasks & async/await.", CYAN),
        ("Day 03", "DOM & Mini-Project", "Event Delegation, Memory optimization, LocalStorage API, Shimmer Skeletons & DevExplorer PRO.", AMBER),
        ("Day 04", "Node.js & FileSystem", "V8 Engine, fs/promises, Path resolution, Process environment & CommonJS vs ES Modules.", GREEN)
    ]

    card_w = 2.75
    gap = 0.24
    start_x = 0.8

    for i, (day, title, desc, col) in enumerate(days_data):
        x = start_x + i * (card_w + gap)
        add_card(s2, x, 1.8, card_w, 4.9, bg_rgb=SURFACE, border_rgb=col)

        tb = s2.shapes.add_textbox(Inches(x + 0.15), Inches(2.0), Inches(card_w - 0.3), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p_day = tf.paragraphs[0]
        p_day.text = day
        p_day.font.size = Pt(14)
        p_day.font.bold = True
        p_day.font.color.rgb = col

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(17)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.space_before = Pt(6)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = MUTED
        p_d.space_before = Pt(14)

    # ==========================================
    # SLIDE 3: DAY 1 - WEB DESIGN & GIT (WITH DIAGRAM IMAGE)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_bg(s3)
    add_header(s3, "Day 01: Modern Web Design Standards & Git Workflows")

    # Left Column: Key Principles
    add_card(s3, 0.8, 1.8, 5.5, 4.9, bg_rgb=SURFACE, border_rgb=INDIGO)
    tb_l3 = s3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_l3 = tb_l3.text_frame
    tf_l3.word_wrap = True

    p_l3_1 = tf_l3.paragraphs[0]
    p_l3_1.text = "🎨 Standards & Git Architecture"
    p_l3_1.font.size = Pt(18)
    p_l3_1.font.bold = True
    p_l3_1.font.color.rgb = CYAN

    bullets_l3 = [
        "Semantic HTML5: Clean hierarchy (<main>, <nav>, <article>) replacing messy <div> soup for optimal SEO & accessibility.",
        "Flexbox (1D) vs Grid (2D): Flexbox for single-axis stacks; CSS Grid for auto-responsive dashboard card matrices.",
        "Tailwind CSS Tokens: Atomic utility classes enforcing standardized spacing, typography, and zero CSS bloat.",
        "Feature-Branching: Isolated branches (feature/task-name) protecting the production 'main' branch.",
        "Conventional Commits: Structured commit logs (feat:, fix:, docs:, refactor:) enabling automated release workflows."
    ]
    for b in bullets_l3:
        p = tf_l3.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11.5)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)

    # Right Column: Visual Diagram Image
    add_card(s3, 6.6, 1.8, 5.9, 4.9, bg_rgb=SURFACE, border_rgb=AMBER)
    git_img_path = os.path.join(assets_dir, "git_branching_diagram.png")
    if os.path.exists(git_img_path):
        s3.shapes.add_picture(git_img_path, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.5))

    # ==========================================
    # SLIDE 4: DAY 2 - JS CORE & FUNCTIONAL PIPELINES
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_bg(s4)
    add_header(s4, "Day 02: Advanced JavaScript Core & Functional Pipelines")

    # Left: ES6+ Syntax
    add_card(s4, 0.8, 1.8, 5.6, 4.9, bg_rgb=SURFACE, border_rgb=CYAN)
    tb_l4 = s4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_l4 = tb_l4.text_frame
    tf_l4.word_wrap = True

    p_l4_1 = tf_l4.paragraphs[0]
    p_l4_1.text = "⚡ ES6+ Scope & Immutability"
    p_l4_1.font.size = Pt(18)
    p_l4_1.font.bold = True
    p_l4_1.font.color.rgb = CYAN

    bullets_l4 = [
        "Block Scope (const / let): Eliminating hoisting bugs and global window pollution.",
        "Destructuring & Rest/Spread: Unpacking API payloads and creating immutable state copies (...state).",
        "Nullish Coalescing (??): Safe default handling for null / undefined values without masking falsy zeros.",
        "Pure Functions: Deterministic logic free of side-effects, laying the foundation for React component state."
    ]
    for b in bullets_l4:
        p = tf_l4.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    # Right: Data Transformations
    add_card(s4, 6.8, 1.8, 5.733, 4.9, bg_rgb=SURFACE, border_rgb=GREEN)
    tb_r4 = s4.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_r4 = tb_r4.text_frame
    tf_r4.word_wrap = True

    p_r4_1 = tf_r4.paragraphs[0]
    p_r4_1.text = "📊 Immutable Data Transformations"
    p_r4_1.font.size = Pt(18)
    p_r4_1.font.bold = True
    p_r4_1.font.color.rgb = GREEN

    bullets_r4 = [
        ".map(): 1-to-1 transformation converting raw GitHub JSON objects into UI presentation cards.",
        ".filter(): Isolating non-forked, active repositories matching user search criteria.",
        ".reduce(): The Swiss Army Knife — single-pass aggregate computations:",
        "   - Total Stars accumulated across all developer repositories",
        "   - Total Forks count measuring ecosystem engagement",
        "   - Primary Language frequency distribution"
    ]
    for b in bullets_r4:
        p = tf_r4.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)

    # ==========================================
    # SLIDE 5: DAY 2 - ASYNC JS & EVENT LOOP (WITH DIAGRAM IMAGE)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_bg(s5)
    add_header(s5, "Day 02 (Deep Dive): The JavaScript Event Loop & Async Architecture")

    # Left: Event Loop Principles
    add_card(s5, 0.8, 1.8, 5.2, 4.9, bg_rgb=SURFACE, border_rgb=INDIGO)
    tb_l5 = s5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.8), Inches(4.5))
    tf_l5 = tb_l5.text_frame
    tf_l5.word_wrap = True

    p_l5_1 = tf_l5.paragraphs[0]
    p_l5_1.text = "🔄 Non-Blocking Execution"
    p_l5_1.font.size = Pt(18)
    p_l5_1.font.bold = True
    p_l5_1.font.color.rgb = INDIGO

    bullets_l5 = [
        "Single-Threaded Call Stack: V8 executes synchronous JavaScript line-by-line (LIFO).",
        "Microtask Queue Priority: Promises (.then(), async/await) execute immediately after Call Stack empties.",
        "Macrotask Queue: setTimeout() and I/O callbacks wait until all microtasks are fully drained.",
        "Concurrent Promise.all(): Fetching User Profile + Repositories in parallel, cutting latency by 50%."
    ]
    for b in bullets_l5:
        p = tf_l5.add_paragraph()
        p.text = "⚡ " + b
        p.font.size = Pt(11.5)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    # Right: High-Res Event Loop Diagram Image
    add_card(s5, 6.3, 1.8, 6.2, 4.9, bg_rgb=SURFACE, border_rgb=AMBER)
    el_img_path = os.path.join(assets_dir, "event_loop_diagram.png")
    if os.path.exists(el_img_path):
        s5.shapes.add_picture(el_img_path, Inches(6.5), Inches(2.0), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 6: DAY 3 - DOM & EVENT DELEGATION (WITH DIAGRAM IMAGE)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_bg(s6)
    add_header(s6, "Day 03: DOM Event Architecture & Memory Optimization")

    # Left: Event Delegation Principles
    add_card(s6, 0.8, 1.8, 5.2, 4.9, bg_rgb=SURFACE, border_rgb=INDIGO)
    tb_l6 = s6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.8), Inches(4.5))
    tf_l6 = tb_l6.text_frame
    tf_l6.word_wrap = True

    p_l6_1 = tf_l6.paragraphs[0]
    p_l6_1.text = "🎯 Event Bubbling & Delegation"
    p_l6_1.font.size = Pt(18)
    p_l6_1.font.bold = True
    p_l6_1.font.color.rgb = CYAN

    bullets_l6 = [
        "Event Bubbling Principle: Click events propagate upwards through the DOM hierarchy.",
        "Single Listener Architecture: 1 listener on the parent container handles 50+ dynamic cards with zero memory leaks.",
        "e.target.closest('.bookmark-btn'): Accurately intercepts clicks even on nested SVGs or text spans.",
        "LocalStorage Sync: JSON.stringify and JSON.parse for instant client-side state persistence across reloads."
    ]
    for b in bullets_l6:
        p = tf_l6.add_paragraph()
        p.text = "✔ " + b
        p.font.size = Pt(11.5)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    # Right: DOM Delegation Diagram Image
    add_card(s6, 6.3, 1.8, 6.2, 4.9, bg_rgb=SURFACE, border_rgb=GREEN)
    dom_img_path = os.path.join(assets_dir, "dom_delegation_diagram.png")
    if os.path.exists(dom_img_path):
        s6.shapes.add_picture(dom_img_path, Inches(6.5), Inches(2.0), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 7: DAY 4 - NODE.JS CORE & FILESYSTEM (WITH DIAGRAM IMAGE)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_bg(s7)
    add_header(s7, "Day 04: Node.js Core, V8 Engine & Asynchronous FileSystem")

    # Left: Node.js Principles
    add_card(s7, 0.8, 1.8, 5.2, 4.9, bg_rgb=SURFACE, border_rgb=GREEN)
    tb_l7 = s7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.8), Inches(4.5))
    tf_l7 = tb_l7.text_frame
    tf_l7.word_wrap = True

    p_l7_1 = tf_l7.paragraphs[0]
    p_l7_1.text = "🟢 Server-Side JavaScript"
    p_l7_1.font.size = Pt(18)
    p_l7_1.font.bold = True
    p_l7_1.font.color.rgb = GREEN

    bullets_l7 = [
        "V8 + Libuv Architecture: JavaScript executing outside browser sandbox with multi-threaded C++ thread pool.",
        "Module Systems: CommonJS (require/exports) vs Modern ES Modules (import/export).",
        "Non-Blocking fs/promises: Asynchronous disk I/O (readFile, writeFile) maintaining 100% server responsiveness.",
        "Safe Cross-Platform Paths: path.join(__dirname, 'data.json') resolving Windows vs POSIX path separators."
    ]
    for b in bullets_l7:
        p = tf_l7.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11.5)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    # Right: Node.js Architecture Diagram Image
    add_card(s7, 6.3, 1.8, 6.2, 4.9, bg_rgb=SURFACE, border_rgb=CYAN)
    node_img_path = os.path.join(assets_dir, "node_arch_diagram.png")
    if os.path.exists(node_img_path):
        s7.shapes.add_picture(node_img_path, Inches(6.5), Inches(2.0), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 8: PROJECT SHOWCASE - DEVEXPLORER PRO V2.0 (WITH LIVE UI SCREENSHOT)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_bg(s8)
    add_header(s8, "Live Demo: DevExplorer PRO v2.0 & Code Walkthrough", "PROJECT SHOWCASE")

    # Left: Features Breakdown
    add_card(s8, 0.8, 1.8, 5.2, 4.9, bg_rgb=SURFACE, border_rgb=AMBER)
    tb_s8 = s8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.8), Inches(4.5))
    tf_s8 = tb_s8.text_frame
    tf_s8.word_wrap = True

    ps8_1 = tf_s8.paragraphs[0]
    ps8_1.text = "🚀 Full-Stack Feature Set"
    ps8_1.font.size = Pt(18)
    ps8_1.font.bold = True
    ps8_1.font.color.rgb = AMBER

    demo_pts = [
        "1. Parallel API Fetching: Promise.all([fetchUser, fetchRepos]) reducing API wait latency by 50%.",
        "2. Live Analytical Dashboard: .reduce() computing Total Stars, Forks, and Top Languages in 1 pass.",
        "3. Event Delegation & Tag Cloud: 1-click exploration of top developers via single container listener.",
        "4. Client-Side Bookmark State: Real-time synchronization with browser LocalStorage API.",
        "5. Resilient UI States: Shimmer loading skeletons, custom 404 cards, and auto-dismissing toasts."
    ]
    for pt in demo_pts:
        p = tf_s8.add_paragraph()
        p.text = pt
        p.font.size = Pt(11.5)
        p.font.color.rgb = WHITE
        p.space_before = Pt(7)

    # Right: Live UI Screenshot Image
    add_card(s8, 6.3, 1.8, 6.2, 4.9, bg_rgb=SURFACE, border_rgb=INDIGO)
    app_img_path = os.path.join(assets_dir, "devexplorer_ui.png")
    if os.path.exists(app_img_path):
        s8.shapes.add_picture(app_img_path, Inches(6.45), Inches(1.95), Inches(5.9), Inches(4.6))

    # ==========================================
    # SLIDE 9: SUMMARY & TRANSITION TO WEEK 2 (REACT)
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_bg(s9)
    add_header(s9, "Summary, Key Takeaways & Readiness for Week 2 (React)", "LOOKING AHEAD")

    # Left: Week 1 Takeaways
    add_card(s9, 0.8, 1.8, 5.7, 4.9, bg_rgb=SURFACE, border_rgb=GREEN)
    tb_l9 = s9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_l9 = tb_l9.text_frame
    tf_l9.word_wrap = True

    p_l9_1 = tf_l9.paragraphs[0]
    p_l9_1.text = "🏆 Week 1 Engineering Mastery"
    p_l9_1.font.size = Pt(18)
    p_l9_1.font.bold = True
    p_l9_1.font.color.rgb = GREEN

    bullets_l9 = [
        "Solid Semantic Foundations: Accessible layouts, CSS Grid systems, and design tokens.",
        "Professional Git Workflows: Protected main branch, atomic conventional commits, and clean PRs.",
        "Asynchronous Mastery: Deep understanding of the Event Loop, Microtasks, and non-blocking I/O.",
        "Memory-Optimized DOM: Event delegation replacing memory-heavy listeners.",
        "Backend Readiness: Server-side JavaScript with Node.js and asynchronous FileSystem."
    ]
    for b in bullets_l9:
        p = tf_l9.add_paragraph()
        p.text = "✔ " + b
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    # Right: Ready for Week 2 (React)
    add_card(s9, 6.8, 1.8, 5.733, 4.9, bg_rgb=SURFACE, border_rgb=CYAN)
    tb_r9 = s9.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_r9 = tb_r9.text_frame
    tf_r9.word_wrap = True

    p_r9_1 = tf_r9.paragraphs[0]
    p_r9_1.text = "⚛️ Ready for Week 2: React Deep Dive"
    p_r9_1.font.size = Pt(18)
    p_r9_1.font.bold = True
    p_r9_1.font.color.rgb = CYAN

    bullets_r9_2 = [
        "Declarative Mindset: Moving from imperative DOM manipulation to Declarative UI State.",
        "JSX & Virtual DOM: Utilizing our ES6+ array methods (.map()) directly in JSX component trees.",
        "State Immutability: Managing component state updates using object and array spread (...state).",
        "Target Deliverable: Interactive Tic-Tac-Toe App with Undo/Redo & Time-Travel history."
    ]
    for b in bullets_r9_2:
        p = tf_r9.add_paragraph()
        p.text = "🚀 " + b
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    out_path = os.path.join(base_dir, "Week_1_Presentation_PRO_Mastery.pptx")
    prs.save(out_path)
    print(f"POWERPOINT GENERATED WITH HIGH-RES DIAGRAMS & UI SCREENSHOTS AT: {out_path}")

if __name__ == "__main__":
    build_presentation()
