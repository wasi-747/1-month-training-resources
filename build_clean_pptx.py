import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_executive_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Executive Modern Dark Palette (Clean, Calm, High-End Corporate)
    BG_DARK = RGBColor(11, 15, 25)         # #0B0F19 Deep Night
    CARD_BG = RGBColor(21, 29, 44)         # #151D2C Sleek Card
    BORDER_SUBTLE = RGBColor(51, 65, 85)   # #334155 Subtle Slate Border
    
    PRIMARY_ACCENT = RGBColor(99, 102, 241)# #6366F1 Modern Indigo
    SKY_BLUE = RGBColor(56, 189, 248)      # #38BDF8 Sky Accent
    TEXT_WHITE = RGBColor(255, 255, 255)   # #FFFFFF
    TEXT_TITLE = RGBColor(241, 245, 249)   # #F1F5F9
    TEXT_BODY = RGBColor(203, 213, 225)    # #CBD5E1
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8

    base_dir = os.path.dirname(os.path.abspath(__file__))
    assets_dir = os.path.join(base_dir, "assets")

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

    def add_header(slide, title_text, category="WEEK 1 TECHNICAL REVIEW"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(8.0), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = SKY_BLUE

        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.6))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

        # Divider
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.733), Inches(0.015))
        div.fill.solid()
        div.fill.fore_color.rgb = BORDER_SUBTLE
        div.line.fill.background()

    def add_card(slide, left, top, width, height):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_SUBTLE
        card.line.width = Pt(1)
        return card

    # ==========================================
    # SLIDE 1: Title Slide (Minimalist & Hero)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1)

    add_card(s1, 0.8, 0.8, 11.733, 5.9)
    tb = s1.shapes.add_textbox(Inches(1.4), Inches(1.3), Inches(10.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "1-MONTH WEB DEVELOPMENT TRAINING • PRESENTATION #1"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = SKY_BLUE
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "Modern Web Foundations &\nAdvanced JavaScript Architecture"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "A comprehensive technical walkthrough covering Semantic Web Standards, Professional Git Workflows, Asynchronous Event Loop, DOM Event Delegation, and Node.js Core File Management."
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(28)

    p3 = tf.add_paragraph()
    p3.text = "🚀 FEATURED LIVE DEMO: DevExplorer PRO v2.0 (GitHub Analytics & Bookmark Manager)"
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = PRIMARY_ACCENT
    p3.space_after = Pt(8)

    p4 = tf.add_paragraph()
    p4.text = "Presenter: Full-Stack Engineering Trainee  |  Date: August 10, 2026"
    p4.font.size = Pt(13)
    p4.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: 4-Day Milestone Roadmap
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2)
    add_header(s2, "Technical Curriculum & Milestone Roadmap", "WEEK 1 OVERVIEW")

    days_data = [
        ("DAY 01", "Web Design & Git", [
            ("HTML5 Semantics", "Replaced <div> soup with <main>, <nav>, <header> for SEO & a11y."),
            ("Flexbox & CSS Grid", "1D component alignment vs 2D responsive grid matrices."),
            ("Professional Git", "Feature-branching (git checkout -b) & Conventional Commits.")
        ]),
        ("DAY 02", "Advanced JS & Async", [
            ("ES6+ Architecture", "Block scope (const/let), destructuring, and spread operator."),
            ("Array Pipelines", ".map(), .filter(), and .reduce() 1-pass analytical totals."),
            ("Event Loop Engine", "Call Stack execution, Microtask Queue vs Macrotasks.")
        ]),
        ("DAY 03", "DOM & Mini-Project", [
            ("Event Delegation", "Single parent listener handles dynamic items via closest()."),
            ("LocalStorage Sync", "Client state persistence using JSON serialization."),
            ("DevExplorer PRO", "Parallel API fetching via Promise.all() + live analytics.")
        ]),
        ("DAY 04", "Node.js & FileSystem", [
            ("V8 & Libuv Engine", "Server-side JavaScript runtime with direct OS access."),
            ("fs/promises Module", "Non-blocking, asynchronous disk read & write operations."),
            ("Path Resolution", "Cross-platform path handling via path.join().")
        ])
    ]

    for idx, (day_tag, day_title, items) in enumerate(days_data):
        x = 0.8 + idx * 2.98
        add_card(s2, x, 1.7, 2.8, 5.2)

        tb = s2.shapes.add_textbox(Inches(x + 0.15), Inches(1.85), Inches(2.5), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_d = tf.paragraphs[0]
        p_d.text = day_tag
        p_d.font.size = Pt(11)
        p_d.font.bold = True
        p_d.font.color.rgb = SKY_BLUE
        p_d.space_after = Pt(4)

        p_t = tf.add_paragraph()
        p_t.text = day_title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.space_after = Pt(16)

        for heading, desc in items:
            p_h = tf.add_paragraph()
            p_h.text = "• " + heading
            p_h.font.size = Pt(12)
            p_h.font.bold = True
            p_h.font.color.rgb = TEXT_TITLE
            p_h.space_after = Pt(2)

            p_desc = tf.add_paragraph()
            p_desc.text = desc
            p_desc.font.size = Pt(10.5)
            p_desc.font.color.rgb = TEXT_MUTED
            p_desc.space_after = Pt(10)

    # ==========================================
    # SLIDE 3: Day 01 - UI Standards & Git
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3)
    add_header(s3, "Modern Web Standards & Team Git Workflows", "DAY 01 ARCHITECTURE")

    # Left: Clean Content Card
    add_card(s3, 0.8, 1.7, 5.7, 5.2)
    tb = s3.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Core Engineering Principles"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

    bullets_s3 = [
        ("Semantic HTML5 Hierarchy", "Structure with <header>, <nav>, <main>, <article>, and <footer> to eliminate unindexed div soup, maximizing SEO visibility and accessibility."),
        ("CSS Layout Engines (Flexbox vs Grid)", "Flexbox for 1D single-axis alignment (navbars, cards) and CSS Grid for 2D auto-responsive dashboard matrices."),
        ("Feature-Branching Strategy", "Production main branch stays locked and stable. Features are isolated in feature/task-name branches and merged via Pull Requests."),
        ("Conventional Commit Standards", "Enforces standardized commit tags (feat:, fix:, docs:, refactor:) for clear, automated changelogs.")
    ]
    for h, d in bullets_s3:
        ph = tf.add_paragraph()
        ph.text = "✔  " + h
        ph.font.size = Pt(12.5)
        ph.font.bold = True
        ph.font.color.rgb = SKY_BLUE
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(12)

    # Right: Diagram
    add_card(s3, 6.8, 1.7, 5.733, 5.2)
    git_img = os.path.join(assets_dir, "git_branching_diagram.png")
    if os.path.exists(git_img):
        s3.shapes.add_picture(git_img, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 4: Day 02 - Advanced JS & Pipelines
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4)
    add_header(s4, "Advanced JavaScript & Functional Data Pipelines", "DAY 02 CORE JS")

    # Left: ES6+ Syntax
    add_card(s4, 0.8, 1.7, 5.7, 5.2)
    tb = s4.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "ES6+ Architecture & Scope"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

    bullets_s4_l = [
        ("Block Scope (const / let)", "Eliminated var hoisting hazards and scope leakage. Using const by default guarantees immutable variable bindings."),
        ("Destructuring & Spread Operator (...)", "Clean object/array unpacking and shallow cloning, forming the foundation for immutable React state updates."),
        ("Nullish Coalescing (??) & Optional Chaining (?.)", "Safely navigates deep nested API objects without throwing uncaught runtime errors.")
    ]
    for h, d in bullets_s4_l:
        ph = tf.add_paragraph()
        ph.text = "⚡  " + h
        ph.font.size = Pt(12.5)
        ph.font.bold = True
        ph.font.color.rgb = SKY_BLUE
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(14)

    # Right: Array Transformations
    add_card(s4, 6.8, 1.7, 5.733, 5.2)
    tb = s4.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.1), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Higher-Order Array Methods"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

    bullets_s4_r = [
        (".map() — 1-to-1 UI Transformation", "Transforms raw backend API payloads into structured UI presentation elements without mutating source arrays."),
        (".filter() — Boolean Data Selection", "Extracts active, non-forked repositories matching user filter criteria."),
        (".reduce() — Single-Pass Analytical Aggregation", "Computes total stars, fork count, and primary language in 1 single pass:\nconst totalStars = repos.reduce((sum, r) => sum + r.stargazers_count, 0);")
    ]
    for h, d in bullets_s4_r:
        ph = tf.add_paragraph()
        ph.text = "📊  " + h
        ph.font.size = Pt(12.5)
        ph.font.bold = True
        ph.font.color.rgb = PRIMARY_ACCENT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(14)

    # ==========================================
    # SLIDE 5: Day 02 Deep Dive - Event Loop
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5)
    add_header(s5, "JavaScript Event Loop & Asynchronous Architecture", "DAY 02 DEEP DIVE")

    add_card(s5, 0.8, 1.7, 5.7, 5.2)
    tb = s5.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Concurrency & Execution Queues"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

    bullets_s5 = [
        ("1. Single-Threaded Call Stack", "V8 executes synchronous functions line-by-line (LIFO). Long I/O tasks are delegated to Web APIs."),
        ("2. Microtask Queue (VIP Priority)", "Promises (.then(), async/await) enter the microtask queue and execute immediately after the Call Stack clears."),
        ("3. Macrotask Queue", "Timer callbacks (setTimeout, setInterval) wait until all pending microtasks are completely drained."),
        ("4. Concurrent Promise.all()", "Dispatches User Profile and Repositories API requests in parallel, cutting network latency by 50%.")
    ]
    for h, d in bullets_s5:
        ph = tf.add_paragraph()
        ph.text = "🔄  " + h
        ph.font.size = Pt(12.5)
        ph.font.bold = True
        ph.font.color.rgb = SKY_BLUE
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    add_card(s5, 6.8, 1.7, 5.733, 5.2)
    loop_img = os.path.join(assets_dir, "event_loop_diagram.png")
    if os.path.exists(loop_img):
        s5.shapes.add_picture(loop_img, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 6: Day 03 - DOM Delegation & LocalStorage
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6)
    add_header(s6, "DOM Event Delegation & State Persistence", "DAY 03 DOM ARCHITECTURE")

    add_card(s6, 0.8, 1.7, 5.7, 5.2)
    tb = s6.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Memory Optimization & Storage"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

    bullets_s6 = [
        ("Event Bubbling Mechanism", "Browser click events propagate upwards through all DOM parent containers automatically."),
        ("Single Container Listener", "Attached 1 single listener to the parent container instead of binding 50+ individual button listeners. Eliminates memory leaks."),
        ("Dynamic Interception via closest()", "event.target.closest('.tag-btn') reliably intercepts clicks on nested child icons and text spans."),
        ("LocalStorage State Persistence", "Serializes bookmarks with JSON.stringify and restores them on page reload with JSON.parse.")
    ]
    for h, d in bullets_s6:
        ph = tf.add_paragraph()
        ph.text = "🎯  " + h
        ph.font.size = Pt(12.5)
        ph.font.bold = True
        ph.font.color.rgb = SKY_BLUE
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    add_card(s6, 6.8, 1.7, 5.733, 5.2)
    dom_img = os.path.join(assets_dir, "dom_delegation_diagram.png")
    if os.path.exists(dom_img):
        s6.shapes.add_picture(dom_img, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 7: Day 04 - Node.js Core Runtime
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7)
    add_header(s7, "Node.js Core Runtime & Asynchronous FileSystem", "DAY 04 BACKEND ARCHITECTURE")

    add_card(s7, 0.8, 1.7, 5.7, 5.2)
    tb = s7.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Server-Side Execution & fs/promises"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

    bullets_s7 = [
        ("Server V8 + Libuv Architecture", "JavaScript executing outside the browser sandbox with multi-threaded C++ Libuv handling non-blocking OS I/O."),
        ("CommonJS Module Standard", "Explicit encapsulation and sharing via require() imports and module.exports."),
        ("Non-Blocking fs/promises Module", "Asynchronous file I/O (readFile, writeFile, appendFile) maintaining full server responsiveness."),
        ("Safe Cross-Platform Paths", "path.join(__dirname, 'file.json') prevents operating system path separator conflicts.")
    ]
    for h, d in bullets_s7:
        ph = tf.add_paragraph()
        ph.text = "🟢  " + h
        ph.font.size = Pt(12.5)
        ph.font.bold = True
        ph.font.color.rgb = SKY_BLUE
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    add_card(s7, 6.8, 1.7, 5.733, 5.2)
    node_img = os.path.join(assets_dir, "node_arch_diagram.png")
    if os.path.exists(node_img):
        s7.shapes.add_picture(node_img, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 8: Project Showcase (DevExplorer PRO)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8)
    add_header(s8, "Live Demonstration: DevExplorer PRO v2.0", "PROJECT SHOWCASE & LIVE DEMO")

    add_card(s8, 0.8, 1.7, 5.7, 5.2)
    tb = s8.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Featured Architecture Highlights"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

    bullets_s8 = [
        ("1. Parallel API Fetching", "Fetches GitHub User Profile and Repositories concurrently via Promise.all() with shimmer loading states."),
        ("2. Real-Time Analytics Engine", "Calculates Total Stars, Fork count, and Top Language using a single-pass .reduce() pipeline."),
        ("3. Event Delegation Tag Cloud", "1-click developer exploration (torvalds, gaearon, wasi-747) handled by a single container listener."),
        ("4. LocalStorage State Sync", "Persists bookmarked developer profiles client-side across browser sessions."),
        ("5. Resilient Error Handling", "Graceful API rate-limit detection, custom 404 cards, and auto-dismissing toast notifications.")
    ]
    for h, d in bullets_s8:
        ph = tf.add_paragraph()
        ph.text = "🚀  " + h
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = PRIMARY_ACCENT
        ph.space_after = Pt(1)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(8)

    add_card(s8, 6.8, 1.7, 5.733, 5.2)
    ui_img = os.path.join(assets_dir, "devexplorer_ui.png")
    if os.path.exists(ui_img):
        s8.shapes.add_picture(ui_img, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 9: Summary & Transition to React
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9)
    add_header(s9, "Summary & Readiness for Week 2 (React Deep Dive)", "LOOKING AHEAD")

    # Left Card
    add_card(s9, 0.8, 1.7, 5.7, 5.2)
    tb = s9.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Week 1 Engineering Mastery"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

    bullets_s9_l = [
        ("Semantic Web Standards", "HTML5 structure, 2D CSS Grid auto-fit matrices, and Tailwind utility tokens."),
        ("Production Git Strategy", "Locked main branch, atomic conventional commits, and Pull Request reviews."),
        ("Asynchronous Architecture", "Call Stack, Microtask queue priority, and non-blocking Promise pipelines."),
        ("Memory-Optimized DOM & Node", "Event delegation, LocalStorage state sync, and asynchronous FileSystem I/O.")
    ]
    for h, d in bullets_s9_l:
        ph = tf.add_paragraph()
        ph.text = "✔  " + h
        ph.font.size = Pt(12.5)
        ph.font.bold = True
        ph.font.color.rgb = SKY_BLUE
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    # Right Card
    add_card(s9, 6.8, 1.7, 5.733, 5.2)
    tb = s9.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.1), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Week 2: React Deep Dive Focus"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

    bullets_s9_r = [
        ("Declarative UI Mindset", "Transitioning from imperative DOM manipulation to Declarative React Component States."),
        ("JSX & Virtual DOM Architecture", "Using ES6+ array pipelines (.map()) directly in JSX component trees."),
        ("State Immutability & Hooks", "Managing state with useState and useEffect cleanly using spread operators (...state)."),
        ("Target Deliverable", "Interactive Tic-Tac-Toe App with Undo/Redo & Time-Travel history.")
    ]
    for h, d in bullets_s9_r:
        ph = tf.add_paragraph()
        ph.text = "🚀  " + h
        ph.font.size = Pt(12.5)
        ph.font.bold = True
        ph.font.color.rgb = PRIMARY_ACCENT
        ph.space_after = Pt(2)

        pd = tf.add_paragraph()
        pd.text = d
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_BODY
        pd.space_after = Pt(10)

    out_file = os.path.join(base_dir, "Week_1_Presentation_Clean_Executive.pptx")
    prs.save(out_file)
    print("SUCCESS: Clean executive presentation saved to", out_file)

if __name__ == "__main__":
    create_executive_deck()
