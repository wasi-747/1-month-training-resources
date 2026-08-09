import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette (Dark Mode Aesthetic)
    BG_COLOR = RGBColor(15, 23, 42)       # #0F172A Deep Slate
    SURFACE_COLOR = RGBColor(30, 41, 59) # #1E293B Card Background
    INDIGO = RGBColor(99, 102, 241)       # #6366F1 Accent Indigo
    AMBER = RGBColor(245, 158, 11)        # #F59E0B Accent Amber
    WHITE = RGBColor(248, 250, 252)       # #F8FAFC Text Main
    MUTED = RGBColor(148, 163, 184)       # #94A3B8 Text Muted
    GREEN = RGBColor(16, 185, 129)        # #10B981 Green

    def apply_bg(slide):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR
        bg_shape.line.fill.background()

    def add_card(slide, left, top, width, height, bg_rgb=SURFACE_COLOR, border_rgb=INDIGO):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_rgb
        card.line.color.rgb = border_rgb
        card.line.width = Pt(1.5)
        return card

    def add_header(slide, title_text, category_text="WEEK 1 PRESENTATION"):
        # Category Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(2.5), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = INDIGO
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.text = category_text
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = WHITE
        p_b.alignment = PP_ALIGN.CENTER

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_bg(slide1)

    # Title Hero Card
    add_card(slide1, 0.8, 0.8, 11.733, 5.9, bg_rgb=SURFACE_COLOR, border_rgb=INDIGO)
    
    # Title Text Box
    tbox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(5.0))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "1-MONTH WEB DEVELOPMENT TRAINING"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = AMBER
    p0.space_after = Pt(10)

    p1 = tf1.add_paragraph()
    p1.text = "Modern Web Foundations, Advanced JS & Node.js"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(15)

    p2 = tf1.add_paragraph()
    p2.text = "Presentation #1 • Comprehensive Week 1 Architecture & Live Code Walkthrough"
    p2.font.size = Pt(16)
    p2.font.color.rgb = MUTED
    p2.space_after = Pt(35)

    p3 = tf1.add_paragraph()
    p3.text = "Presenter: Full-Stack Engineering Trainee"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = INDIGO

    # ==========================================
    # SLIDE 2: Modern Web Design Standards
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide2)
    add_header(slide2, "1. Modern Web Design & Layout Architecture")

    # Card 1: Semantic HTML5
    add_card(slide2, 0.8, 1.8, 3.6, 5.0)
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Semantic HTML5"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AMBER
    p.space_after = Pt(14)

    bullets1 = [
        "Replaces <div> soup with meaningful tags (<main>, <header>, <nav>, <footer>).",
        "Maximizes Search Engine Optimization (SEO) indexing.",
        "Ensures web accessibility (a11y) for screen readers."
    ]
    for b in bullets1:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)

    # Card 2: Layout Engines
    add_card(slide2, 4.8, 1.8, 3.6, 5.0)
    tb = slide2.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Flexbox vs. Grid"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.space_after = Pt(14)

    bullets2 = [
        "Flexbox (1D): Single-axis alignment for navbars, row stacks, and centering.",
        "CSS Grid (2D): Multi-column layouts for dashboard card grids.",
        "Responsive media queries for mobile-first layouts."
    ]
    for b in bullets2:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)

    # Card 3: Tailwind CSS
    add_card(slide2, 8.8, 1.8, 3.733, 5.0)
    tb = slide2.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Tailwind CSS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(14)

    bullets3 = [
        "Utility-first CSS framework for rapid UI development.",
        "Standardized design tokens (padding, spacing, HSL color palettes).",
        "Zero custom CSS overhead with reusable atomic classes."
    ]
    for b in bullets3:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 3: Git & GitHub Team Workflows
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_bg(slide3)
    add_header(slide3, "2. Professional Git & GitHub Team Workflows")

    add_card(slide3, 0.8, 1.8, 5.6, 5.0)
    tb = slide3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Feature-Branching Strategy"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.space_after = Pt(14)

    bullets_git1 = [
        "Production Protection: The 'main' branch remains locked and stable.",
        "Isolated Development: Create feature branches via git checkout -b feature/task-name.",
        "Prevents merge conflicts across engineering teams.",
        "Enforces Pull Request (PR) code reviews before merging."
    ]
    for b in bullets_git1:
        p = tf.add_paragraph()
        p.text = "✔ " + b
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
        p.space_after = Pt(12)

    add_card(slide3, 6.8, 1.8, 5.733, 5.0)
    tb = slide3.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Conventional Commit Standards"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AMBER
    p.space_after = Pt(14)

    bullets_git2 = [
        "feat: Adding new functional features (e.g. feat: add bookmarking).",
        "fix: Patching software bugs (e.g. fix: resolve flexbox alignment).",
        "docs: Updating documentation and README study guides.",
        "style: Formatting or CSS design adjustments without logic changes."
    ]
    for b in bullets_git2:
        p = tf.add_paragraph()
        p.text = "📌 " + b
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
        p.space_after = Pt(12)

    # ==========================================
    # SLIDE 4: Advanced JavaScript & Async Execution
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_bg(slide4)
    add_header(slide4, "3. Advanced JavaScript & Asynchronous Architecture")

    # Left: JS Core & Array Methods
    add_card(slide4, 0.8, 1.8, 5.6, 5.0)
    tb = slide4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ES6+ Data Transformations"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AMBER
    p.space_after = Pt(14)

    js_bullets1 = [
        "Block Scope: const for immutable bindings, let for mutable variables.",
        "Arrow Functions: Clean syntax with implicit return capability.",
        ".map(): Transforms elements into a new array of equal length.",
        ".filter(): Extracts elements matching boolean predicates.",
        ".reduce(): Accumulates array values into single analytical totals."
    ]
    for b in js_bullets1:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)

    # Right: Async & Event Loop
    add_card(slide4, 6.8, 1.8, 5.733, 5.0)
    tb = slide4.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Event Loop & Promises"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.space_after = Pt(14)

    js_bullets2 = [
        "Single-Threaded Event Loop: Handles non-blocking execution via Call Stack and Microtask Queue.",
        "Microtask Priority: Promises execute immediately after Call Stack clears.",
        "async / await: Clean, readable asynchronous control flow syntax.",
        "try / catch: Error boundaries preventing unhandled promise rejections."
    ]
    for b in js_bullets2:
        p = tf.add_paragraph()
        p.text = "⚡ " + b
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 5: DOM Event Architecture & Node.js
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_bg(slide5)
    add_header(slide5, "4. DOM Event Architecture & Node.js Core")

    # Left: DOM & LocalStorage
    add_card(slide5, 0.8, 1.8, 5.6, 5.0)
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DOM & LocalStorage API"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(14)

    dom_bullets = [
        "Event Delegation: Single listener on parent container using event.target.closest().",
        "Memory Optimization: Eliminates hundreds of individual DOM event handlers.",
        "LocalStorage API: Client-side state persistence (JSON.stringify / JSON.parse).",
        "UI State Feedback: Loading spinners, Error cards, and Toast notifications."
    ]
    for b in dom_bullets:
        p = tf.add_paragraph()
        p.text = "🎯 " + b
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)

    # Right: Node.js fs/promises
    add_card(slide5, 6.8, 1.8, 5.733, 5.0)
    tb = slide5.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Node.js Core & File System"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.space_after = Pt(14)

    node_bullets = [
        "V8 Execution Engine: Running JavaScript on backend servers outside the browser.",
        "fs/promises Module: Non-blocking asynchronous file operations (readFile, writeFile).",
        "Global Context: process.env for configuration and __dirname for path resolution.",
        "HTTP Server Basics: Built-in http module for RESTful response handling."
    ]
    for b in node_bullets:
        p = tf.add_paragraph()
        p.text = "🟢 " + b
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 6: Live Code Demo & Conclusion
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_bg(slide6)
    add_header(slide6, "5. Live Code Demo — DevExplorer PRO v2.0", "PROJECT DEMONSTRATION")

    add_card(slide6, 0.8, 1.8, 11.733, 5.0, bg_rgb=SURFACE_COLOR, border_rgb=AMBER)
    tb = slide6.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "Live Demonstration: GitHub Explorer & Bookmark Manager"
    p0.font.size = Pt(22)
    p0.font.bold = True
    p0.font.color.rgb = WHITE
    p0.space_after = Pt(15)

    demo_bullets = [
        "Parallel API Requests: Fetching user profiles and repositories concurrently via Promise.all().",
        "Event Delegation: Single container listener handling quick developer tags and bookmarking.",
        "Analytics Engine: Live computation of Total Stars, Forks, and Top Language via .reduce().",
        "LocalStorage Persistence: Saving favorite developer profiles across page reloads.",
        "GitHub Repository: https://github.com/wasi-747/1-month-training-resources"
    ]
    for b in demo_bullets:
        p = tf.add_paragraph()
        p.text = "🚀 " + b
        p.font.size = Pt(15)
        p.font.color.rgb = MUTED
        p.space_after = Pt(12)

    # Output file path
    out_path = os.path.join("d:\\Study\\Projects\\1 month training resources", "Presentation_1_Modern_Web_Foundations.pptx")
    prs.save(out_path)
    print("SUCCESS:", out_path)

if __name__ == "__main__":
    create_deck()
