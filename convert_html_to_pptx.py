import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Tokens from HTML Deck
    BG_SLATE = RGBColor(8, 12, 22)         # #080C16
    CARD_BG = RGBColor(17, 24, 39)         # #111827
    CARD_INNER = RGBColor(22, 32, 50)      # #162032
    TOPIC_BG = RGBColor(11, 15, 25)        # rgba(11, 15, 25, 0.6)
    
    INDIGO = RGBColor(99, 102, 241)        # #6366F1
    INDIGO_LIGHT = RGBColor(165, 180, 252) # #A5B4FC
    CYAN = RGBColor(6, 182, 212)           # #06B6D4
    CYAN_LIGHT = RGBColor(103, 232, 249)   # #67E8F9
    AMBER = RGBColor(245, 158, 11)         # #F59E0B
    AMBER_LIGHT = RGBColor(252, 211, 77)   # #FCD34D
    EMERALD = RGBColor(16, 185, 129)       # #10B981
    EMERALD_LIGHT = RGBColor(110, 231, 183)# #6EE7B7
    WHITE = RGBColor(248, 250, 252)        # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8
    TEXT_BODY = RGBColor(203, 213, 225)    # #CBD5E1

    base_dir = os.path.dirname(os.path.abspath(__file__))
    assets_dir = os.path.join(base_dir, "assets")

    def apply_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_SLATE
        bg.line.fill.background()

    def add_header(slide, title_text, badge_text, badge_color=INDIGO, badge_light=INDIGO_LIGHT):
        # Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(3.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = CARD_INNER
        badge.line.color.rgb = badge_color
        badge.line.width = Pt(1.2)
        tf_b = badge.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.text = badge_text.upper()
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = badge_light
        p_b.alignment = PP_ALIGN.CENTER

        # Title
        tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.6))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Divider line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(40, 50, 70)
        line.line.fill.background()

    def add_card(slide, left, top, width, height, border_color=INDIGO, bg_color=CARD_INNER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    # ==========================================
    # SLIDE 1: Title & Hero
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1)
    
    # Outer Hero Card
    add_card(s1, 0.8, 0.7, 11.733, 6.1, border_color=INDIGO, bg_color=CARD_BG)
    
    # Badges row
    b1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.1), Inches(4.2), Inches(0.35))
    b1.fill.solid()
    b1.fill.fore_color.rgb = CARD_INNER
    b1.line.color.rgb = INDIGO
    b1.line.width = Pt(1)
    tf = b1.text_frame
    p = tf.paragraphs[0]
    p.text = "1-MONTH WEB DEVELOPMENT TRAINING • PRESENTATION #1"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = INDIGO_LIGHT
    p.alignment = PP_ALIGN.CENTER

    b2 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.6), Inches(1.1), Inches(2.8), Inches(0.35))
    b2.fill.solid()
    b2.fill.fore_color.rgb = CARD_INNER
    b2.line.color.rgb = CYAN
    b2.line.width = Pt(1)
    tf = b2.text_frame
    p = tf.paragraphs[0]
    p.text = "COMPREHENSIVE WEEK 1 REVIEW"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CYAN_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Main Hero Title
    tbox = s1.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11.0), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "Modern Web Foundations &\nAdvanced JavaScript Architecture"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = "Full Architecture Breakdown: Semantic Layouts, Git Team Workflows, Async Event Loop, DOM Delegation & Node.js Core."
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT_MUTED

    # Day badges
    days = [
        ("Day 1: UI & Git", CYAN, CYAN_LIGHT),
        ("Day 2: Async JS & Event Loop", AMBER, AMBER_LIGHT),
        ("Day 3: DOM & DevExplorer PRO", EMERALD, EMERALD_LIGHT),
        ("Day 4: Node.js & FS", INDIGO, INDIGO_LIGHT)
    ]
    cur_x = 1.2
    for d_text, b_col, b_lcol in days:
        db = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cur_x), Inches(4.3), Inches(2.6), Inches(0.38))
        db.fill.solid()
        db.fill.fore_color.rgb = CARD_INNER
        db.line.color.rgb = b_col
        db.line.width = Pt(1)
        tf = db.text_frame
        p = tf.paragraphs[0]
        p.text = d_text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = b_lcol
        p.alignment = PP_ALIGN.CENTER
        cur_x += 2.8

    # Live demo callout box
    demo_box = add_card(s1, 1.2, 5.0, 10.933, 1.3, border_color=AMBER, bg_color=CARD_INNER)
    tb_demo = s1.shapes.add_textbox(Inches(1.4), Inches(5.1), Inches(10.5), Inches(1.1))
    tf = tb_demo.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚀 FEATURED LIVE PROJECT DEMO: DevExplorer PRO v2.0 (GitHub Analytics & Bookmark Manager)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = AMBER_LIGHT
    p.space_after = Pt(6)

    p_sub = tf.add_paragraph()
    p_sub.text = "Parallel API Fetching via Promise.all() • Single-Pass .reduce() Statistics • LocalStorage JSON State Sync"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = CYAN_LIGHT

    # ==========================================
    # SLIDE 2: 4-Day Roadmap
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_header(s2, "Week 1 Technical Curriculum & Milestone Roadmap", "4 DAYS • 4 MILESTONES", CYAN, CYAN_LIGHT)

    col_data = [
        ("Day 01", "Web Design & Git", INDIGO, INDIGO_LIGHT, [
            ("HTML5 Semantics", "Replaced div soup with <main>, <nav>, <article> for SEO & a11y."),
            ("Flexbox & Grid", "1D single-axis stacks vs 2D responsive matrix layouts."),
            ("Professional Git", "Feature-branching & Conventional Commits (feat:, fix:).")
        ], "✔ Clean UI & Branching"),
        ("Day 02", "Advanced JS & Async", CYAN, CYAN_LIGHT, [
            ("ES6+ Immutability", "Block scope, destructuring, spread operator, nullish coalescing."),
            ("Functional Pipelines", ".map(), .filter(), and .reduce() 1-pass aggregate calculations."),
            ("Event Loop & Tasks", "Call Stack, Microtask Queue (VIP Promises) vs Macrotasks.")
        ], "⚡ Event Loop Mastery"),
        ("Day 03", "DOM & Mini-Project", AMBER, AMBER_LIGHT, [
            ("Event Delegation", "1 single parent listener handles 50+ cards via closest()."),
            ("LocalStorage API", "Client state sync using JSON.stringify / parse."),
            ("DevExplorer PRO", "Parallel API fetching via Promise.all() + skeleton loaders.")
        ], "🚀 Live Application"),
        ("Day 04", "Node.js & FileSystem", EMERALD, EMERALD_LIGHT, [
            ("Server-Side V8", "JavaScript outside browser with direct OS and file system access."),
            ("fs/promises Module", "Asynchronous, non-blocking disk I/O maintaining server speed."),
            ("Path Resolution", "Cross-platform path handling via path.join().")
        ], "🟢 Backend Ready")
    ]

    for idx, (day_lbl, title, col_border, col_light, topics, pill_txt) in enumerate(col_data):
        x = 0.8 + idx * 2.98
        add_card(s2, x, 1.7, 2.8, 5.3, border_color=col_border)
        
        # Day Badge
        db = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(1.85), Inches(1.2), Inches(0.28))
        db.fill.solid()
        db.fill.fore_color.rgb = CARD_BG
        db.line.color.rgb = col_border
        db.line.width = Pt(1)
        tf = db.text_frame
        p = tf.paragraphs[0]
        p.text = day_lbl
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = col_light
        p.alignment = PP_ALIGN.CENTER

        # Title
        tb = s2.shapes.add_textbox(Inches(x + 0.15), Inches(2.15), Inches(2.5), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col_light

        # Topics
        top_y = 2.65
        for t_label, t_desc in topics:
            tb_t = s2.shapes.add_textbox(Inches(x + 0.15), Inches(top_y), Inches(2.5), Inches(1.1))
            tf_t = tb_t.text_frame
            tf_t.word_wrap = True
            p = tf_t.paragraphs[0]
            p.text = "• " + t_label
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = CYAN_LIGHT
            p.space_after = Pt(2)

            p_d = tf_t.add_paragraph()
            p_d.text = t_desc
            p_d.font.size = Pt(9.5)
            p_d.font.color.rgb = TEXT_BODY
            top_y += 1.05

        # Callout Pill
        cp = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(6.45), Inches(2.5), Inches(0.35))
        cp.fill.solid()
        cp.fill.fore_color.rgb = CARD_BG
        cp.line.color.rgb = col_border
        cp.line.width = Pt(1)
        tf = cp.text_frame
        p = tf.paragraphs[0]
        p.text = pill_txt
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = col_light
        p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 3: Day 01 - UI Standards & Git (Diagram)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_header(s3, "Day 01: Modern Web Design Standards & Professional Git", "UI LAYOUTS & VERSION CONTROL", INDIGO, INDIGO_LIGHT)

    # Left Card
    add_card(s3, 0.8, 1.7, 5.7, 5.3, border_color=INDIGO)
    tb = s3.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎨 Architecture Standards & Git Strategy"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = INDIGO_LIGHT
    p.space_after = Pt(12)

    topics_s3 = [
        ("Semantic HTML5 Hierarchy", "Structure with <header>, <nav>, <main>, <article>, and <footer> instead of unindexed <div> soup to maximize SEO and accessibility."),
        ("CSS Layout Engines (Flexbox vs Grid)", "Flexbox for 1D single-axis alignment (navbars, toolbars) and CSS Grid for 2D auto-responsive matrices (repeat(auto-fit, minmax(280px, 1fr)))."),
        ("Feature-Branching & Conventional Commits", "Production main branch stays locked. Features developed in feature/task-name and merged via Pull Request with standard commit prefixes (feat:, fix:, docs:).")
    ]
    for lbl, desc in topics_s3:
        p_l = tf.add_paragraph()
        p_l.text = "📌 " + lbl
        p_l.font.size = Pt(13)
        p_l.font.bold = True
        p_l.font.color.rgb = CYAN_LIGHT
        p_l.space_after = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_BODY
        p_d.space_after = Pt(10)

    # Right Image Card
    add_card(s3, 6.8, 1.7, 5.733, 5.3, border_color=INDIGO, bg_color=CARD_BG)
    git_img_path = os.path.join(assets_dir, "git_branching_diagram.png")
    if os.path.exists(git_img_path):
        s3.shapes.add_picture(git_img_path, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 4: Day 02 - Advanced JS & Pipelines
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_header(s4, "Day 02: Advanced JavaScript Core & Functional Data Pipelines", "ES6+ IMMUTABILITY & ARRAY METHODS", CYAN, CYAN_LIGHT)

    # Left Card
    add_card(s4, 0.8, 1.7, 5.7, 5.3, border_color=CYAN)
    tb = s4.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚡ ES6+ Syntax & Scope Integrity"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = CYAN_LIGHT
    p.space_after = Pt(12)

    topics_s4_left = [
        ("Block Scoping (const / let)", "Eliminated var hoisting hazards and global window pollution; immutable reference bindings."),
        ("Destructuring & Spread Operator (...state)", "Clean payload unpacking and non-destructive object/array cloning, preparing for React state immutability."),
        ("Nullish Coalescing (??) & Optional Chaining (?.)", "Safe nested property traversal without throwing runtime errors on null or undefined.")
    ]
    for lbl, desc in topics_s4_left:
        p_l = tf.add_paragraph()
        p_l.text = "⚡ " + lbl
        p_l.font.size = Pt(13)
        p_l.font.bold = True
        p_l.font.color.rgb = CYAN_LIGHT
        p_l.space_after = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_BODY
        p_d.space_after = Pt(10)

    # Right Card
    add_card(s4, 6.8, 1.7, 5.733, 5.3, border_color=EMERALD)
    tb = s4.shapes.add_textbox(Inches(7.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📊 Immutable Data Transformations"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = EMERALD_LIGHT
    p.space_after = Pt(12)

    topics_s4_right = [
        (".map() — 1-to-1 UI Transformation", "Transforms raw GitHub repository objects into structured DOM presentation cards."),
        (".filter() — Data Selection", "Extracts active, non-forked public repositories matching user filter criteria."),
        (".reduce() — Single-Pass Analytical Totals", "Computes Total Stars, Fork count, and Language distribution simultaneously in 1 pass:\nconst totalStars = repos.reduce((sum, r) => sum + r.stargazers_count, 0);")
    ]
    for lbl, desc in topics_s4_right:
        p_l = tf.add_paragraph()
        p_l.text = "📊 " + lbl
        p_l.font.size = Pt(13)
        p_l.font.bold = True
        p_l.font.color.rgb = EMERALD_LIGHT
        p_l.space_after = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_BODY
        p_d.space_after = Pt(10)

    # ==========================================
    # SLIDE 5: Day 02 Deep Dive - Event Loop (Diagram)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_header(s5, "Day 02 (Deep Dive): The JavaScript Event Loop & Async Engine", "NON-BLOCKING EXECUTION ARCHITECTURE", AMBER, AMBER_LIGHT)

    add_card(s5, 0.8, 1.7, 5.7, 5.3, border_color=AMBER)
    tb = s5.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔄 Execution Stack & Priority Queues"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = AMBER_LIGHT
    p.space_after = Pt(12)

    topics_s5 = [
        ("1. Single-Threaded Call Stack", "V8 engine executes synchronous code line-by-line (LIFO). Long I/O tasks are delegated to Web APIs."),
        ("2. Microtask Queue (VIP Priority)", "Promises (.then(), async/await) enter the microtask queue and run immediately after Call Stack clears."),
        ("3. Macrotask Queue", "setTimeout and I/O events wait until all pending microtasks are completely drained."),
        ("4. Concurrent Promise.all()", "Fetches User Profile + Repositories in parallel, slashing network latency by 50%.")
    ]
    for lbl, desc in topics_s5:
        p_l = tf.add_paragraph()
        p_l.text = "⚡ " + lbl
        p_l.font.size = Pt(12.5)
        p_l.font.bold = True
        p_l.font.color.rgb = AMBER_LIGHT
        p_l.space_after = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_BODY
        p_d.space_after = Pt(8)

    add_card(s5, 6.8, 1.7, 5.733, 5.3, border_color=AMBER, bg_color=CARD_BG)
    loop_img_path = os.path.join(assets_dir, "event_loop_diagram.png")
    if os.path.exists(loop_img_path):
        s5.shapes.add_picture(loop_img_path, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 6: Day 03 - DOM Delegation (Diagram)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_header(s6, "Day 03: DOM Event Architecture & Memory Optimization", "EVENT BUBBLING & LOCALSTORAGE SYNC", EMERALD, EMERALD_LIGHT)

    add_card(s6, 0.8, 1.7, 5.7, 5.3, border_color=EMERALD)
    tb = s6.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 Event Delegation & State Persistence"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = EMERALD_LIGHT
    p.space_after = Pt(12)

    topics_s6 = [
        ("Event Bubbling Principle", "Browser click events automatically propagate upwards from target child elements to parent containers."),
        ("Single Listener Architecture", "Attached 1 single event listener on the parent grid instead of 50+ individual button listeners. Eliminates memory leaks."),
        ("Dynamic Element Interception", "event.target.closest('.bookmark-btn') reliably captures clicks even on nested SVG icons or text spans inside dynamic cards."),
        ("LocalStorage Persistence API", "Client state serialized with JSON.stringify and retrieved via JSON.parse across page reloads.")
    ]
    for lbl, desc in topics_s6:
        p_l = tf.add_paragraph()
        p_l.text = "🎯 " + lbl
        p_l.font.size = Pt(12.5)
        p_l.font.bold = True
        p_l.font.color.rgb = EMERALD_LIGHT
        p_l.space_after = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_BODY
        p_d.space_after = Pt(8)

    add_card(s6, 6.8, 1.7, 5.733, 5.3, border_color=EMERALD, bg_color=CARD_BG)
    dom_img_path = os.path.join(assets_dir, "dom_delegation_diagram.png")
    if os.path.exists(dom_img_path):
        s6.shapes.add_picture(dom_img_path, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 7: Day 04 - Node.js Architecture (Diagram)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_header(s7, "Day 04: Node.js Core, V8 Engine & Asynchronous FileSystem", "BACKEND RUNTIME ARCHITECTURE", CYAN, CYAN_LIGHT)

    add_card(s7, 0.8, 1.7, 5.7, 5.3, border_color=CYAN)
    tb = s7.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🟢 Server-Side JavaScript & fs/promises"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = CYAN_LIGHT
    p.space_after = Pt(12)

    topics_s7 = [
        ("Server-Side V8 + Libuv Architecture", "JavaScript executing outside browser sandbox with multi-threaded C++ Libuv thread pool handling non-blocking OS I/O."),
        ("Module Systems (CommonJS vs ES Modules)", "Standardized on require / module.exports vs modern ECMAScript import / export modules."),
        ("Non-Blocking fs/promises Module", "Asynchronous disk I/O (readFile, writeFile) maintaining 100% server responsiveness without thread blocking."),
        ("Safe Cross-Platform Paths", "path.join(__dirname, 'data.json') preventing OS separator conflicts (Windows \\ vs POSIX /).")
    ]
    for lbl, desc in topics_s7:
        p_l = tf.add_paragraph()
        p_l.text = "🟢 " + lbl
        p_l.font.size = Pt(12.5)
        p_l.font.bold = True
        p_l.font.color.rgb = CYAN_LIGHT
        p_l.space_after = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_BODY
        p_d.space_after = Pt(8)

    add_card(s7, 6.8, 1.7, 5.733, 5.3, border_color=CYAN, bg_color=CARD_BG)
    node_img_path = os.path.join(assets_dir, "node_arch_diagram.png")
    if os.path.exists(node_img_path):
        s7.shapes.add_picture(node_img_path, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 8: Live Demonstration (Screenshot)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_header(s8, "Live Demonstration: DevExplorer PRO v2.0", "PROJECT SHOWCASE & LIVE WALKTHROUGH", AMBER, AMBER_LIGHT)

    add_card(s8, 0.8, 1.7, 5.7, 5.3, border_color=AMBER)
    tb = s8.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚀 Core Feature Architecture"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = AMBER_LIGHT
    p.space_after = Pt(12)

    topics_s8 = [
        ("1. Parallel API Requests (Promise.all)", "User profile and repository payloads fetched concurrently with shimmer skeleton feedback."),
        ("2. Real-Time Analytics Engine (.reduce)", "Live calculation of Total Stars, Fork count, and dominant programming language in 1 single pass."),
        ("3. Event Delegation & Quick Tag Cloud", "1-click exploration of top developers (torvalds, gaearon, wasi-747) via single container listener."),
        ("4. LocalStorage Bookmark Persistence", "Saved developer profiles stored client-side and dynamically rendered on page load."),
        ("5. Resilient UI States & Error Boundaries", "Graceful rate-limit detection, custom 404 cards, and auto-dismissing toast notifications.")
    ]
    for lbl, desc in topics_s8:
        p_l = tf.add_paragraph()
        p_l.text = "🚀 " + lbl
        p_l.font.size = Pt(12)
        p_l.font.bold = True
        p_l.font.color.rgb = AMBER_LIGHT
        p_l.space_after = Pt(1)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_BODY
        p_d.space_after = Pt(6)

    add_card(s8, 6.8, 1.7, 5.733, 5.3, border_color=AMBER, bg_color=CARD_BG)
    ui_img_path = os.path.join(assets_dir, "devexplorer_ui.png")
    if os.path.exists(ui_img_path):
        s8.shapes.add_picture(ui_img_path, Inches(7.0), Inches(1.9), width=Inches(5.333))

    # ==========================================
    # SLIDE 9: Summary & Week 2 Readiness
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_header(s9, "Summary & Readiness for Week 2 (React Deep Dive)", "LOOKING AHEAD", EMERALD, EMERALD_LIGHT)

    # Left: Week 1 Recap
    add_card(s9, 0.8, 1.7, 5.7, 5.3, border_color=EMERALD)
    tb = s9.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏆 Week 1 Engineering Mastery"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = EMERALD_LIGHT
    p.space_after = Pt(12)

    topics_s9_left = [
        ("Semantic & Responsive Web Design", "Mastered HTML5 hierarchy, CSS Grid auto-fit matrices, and Tailwind utility tokens."),
        ("Professional Git Workflows", "Protected main branch, atomic conventional commits, and clean Pull Request reviews."),
        ("Asynchronous & Event Loop Mastery", "Deep comprehension of Call Stack, Microtasks vs Macrotasks, and parallel Promise pipelines."),
        ("Memory-Optimized DOM & Node.js Core", "Event delegation, LocalStorage state sync, and asynchronous FileSystem operations.")
    ]
    for lbl, desc in topics_s9_left:
        p_l = tf.add_paragraph()
        p_l.text = "✔ " + lbl
        p_l.font.size = Pt(12.5)
        p_l.font.bold = True
        p_l.font.color.rgb = EMERALD_LIGHT
        p_l.space_after = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_BODY
        p_d.space_after = Pt(8)

    # Right: Week 2 Preview
    add_card(s9, 6.8, 1.7, 5.733, 5.3, border_color=CYAN)
    tb = s9.shapes.add_textbox(Inches(7.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚛️ Ready for Week 2: React Deep Dive"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = CYAN_LIGHT
    p.space_after = Pt(12)

    topics_s9_right = [
        ("Declarative UI Mindset", "Transitioning from imperative DOM manipulation to Declarative React Component States."),
        ("JSX & Virtual DOM Architecture", "Utilizing our ES6+ array pipelines (.map()) directly in JSX component trees."),
        ("State Synchronization & Immutability", "Managing state updates cleanly using spread operators (...state) without mutating original references."),
        ("Target Deliverable", "Interactive Tic-Tac-Toe App with Undo/Redo & Time-Travel history.")
    ]
    for lbl, desc in topics_s9_right:
        p_l = tf.add_paragraph()
        p_l.text = "🚀 " + lbl
        p_l.font.size = Pt(12.5)
        p_l.font.bold = True
        p_l.font.color.rgb = CYAN_LIGHT
        p_l.space_after = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_BODY
        p_d.space_after = Pt(8)

    # Save Output Presentation
    out_file = os.path.join(base_dir, "Week_1_Interactive_Presentation.pptx")
    prs.save(out_file)
    print("SUCCESS: Saved presentation to", out_file)

if __name__ == "__main__":
    build_presentation()
